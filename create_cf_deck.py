"""
ステップアップナビ クラウドファンディング（Readyfor）向け資料
全14枚 — ダークネイビー × エメラルドグリーン デザイン
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── カラーパレット ─────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)  # #0f172a メイン背景
BG_CARD      = RGBColor(0x1E, 0x29, 0x3B)  # #1e293b カード背景
EMERALD      = RGBColor(0x10, 0xB9, 0x81)  # #10b981 アクセント
TEXT_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)  # #f8fafc 白テキスト
TEXT_GRAY    = RGBColor(0xCB, 0xD5, 0xE1)  # #cbd5e1 グレーテキスト
TEXT_MUTED   = RGBColor(0x94, 0xA3, 0xB8)  # #94a3b8 薄いグレー
CARD_BORDER  = RGBColor(0x1E, 0x40, 0x2F)  # グリーン薄いボーダー
EMERALD_DARK = RGBColor(0x10, 0x4A, 0x38)  # ダークエメラルド（ヘッダー・ボーダー用）
EMERALD_MID  = RGBColor(0x10, 0x7A, 0x55)  # ミッドエメラルド

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]


# ─── ヘルパー関数 ────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill_color, transparency=0.0):
    """楕円を追加（透明度オプション: 0.0=不透明, 1.0=完全透明）"""
    shape = slide.shapes.add_shape(9, left, top, width, height)  # 9 = OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if transparency > 0:
        srgbClr = shape._element.find('.//' + qn('a:srgbClr'))
        if srgbClr is not None:
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', str(int((1 - transparency) * 100000)))
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT_WHITE,
                align=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=16, bold=False, color=TEXT_WHITE,
                          align=PP_ALIGN.LEFT, font_name="Meiryo",
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def add_deco_circle(slide):
    """右上に薄いエメラルドの装飾円（10%不透明）"""
    add_oval(slide,
             left=Inches(11.0), top=Inches(-1.0),
             width=Inches(3.0), height=Inches(3.0),
             fill_color=EMERALD, transparency=0.90)


def add_footer(slide):
    """右下フッター"""
    add_textbox(slide, "🌱 ステップアップナビ",
                Inches(10.5), Inches(7.1), Inches(2.7), Inches(0.35),
                font_size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def body_slide_base(slide, title_text, subtitle_text=None):
    """本文スライド共通: ダークネイビー背景 + 左縦グリーンボーダー + タイトル"""
    add_rect(slide, 0, 0, W, H, fill_color=BG_DARK)
    add_deco_circle(slide)
    # タイトル左の太い縦ボーダー
    add_rect(slide, Inches(0.3), Inches(0.22), Inches(0.08), Inches(0.52), fill_color=EMERALD)
    add_textbox(slide, title_text,
                Inches(0.48), Inches(0.22), Inches(10), Inches(0.6),
                font_size=28, bold=True, color=TEXT_WHITE)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.48), Inches(0.85), Inches(11), Inches(0.4),
                    font_size=16, color=EMERALD)
    add_footer(slide)


def add_card(slide, left, top, width, height):
    """カードコンポーネント（BG_CARD背景 + エメラルド薄ボーダー）"""
    return add_rect(slide, left, top, width, height,
                    fill_color=BG_CARD, line_color=EMERALD_DARK, line_width=Pt(1.2))


def add_badge(slide, text, left, top, width, height, font_size=12):
    """バッジコンポーネント（EMERALD背景 + BG_DARK文字）"""
    add_rect(slide, left, top, width, height, fill_color=EMERALD)
    add_textbox(slide, text,
                left + Inches(0.05), top, width - Inches(0.1), height,
                font_size=font_size, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


def add_dark_table(slide, headers, rows, left, top, width, height, font_size=13):
    """ダークテーブル（EMERALD ヘッダー + 交互行）"""
    col_count = len(headers)
    row_count = len(rows) + 1
    tbl = slide.shapes.add_table(row_count, col_count, left, top, width, height).table

    col_widths = [int(width * r) for r in [0.18, 0.82]]
    if col_count == 2:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    else:
        col_w = width // col_count
        for i in range(col_count):
            tbl.columns[i].width = col_w

    def style_cell(cell, text, bg, fg=TEXT_WHITE, bold=False, sz=font_size,
                   align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Meiryo"

    for ci, h in enumerate(headers):
        style_cell(tbl.cell(0, ci), h, EMERALD, BG_DARK, bold=True)

    for ri, row in enumerate(rows):
        bg = BG_CARD if ri % 2 == 0 else BG_DARK
        for ci, val in enumerate(row):
            align = PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER
            style_cell(tbl.cell(ri + 1, ci), val, bg, TEXT_GRAY, align=align)

    return tbl


# ══════════════════════════════════════════════════════════════
# スライド 1：表紙
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK)
add_rect(slide1, 0, 0, W, H, fill_color=BG_DARK)
add_deco_circle(slide1)

# 絵文字「🌱」大きく中央上部
add_textbox(slide1, "🌱",
            Inches(5.5), Inches(0.65), Inches(2.33), Inches(1.1),
            font_size=64, bold=False, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# タイトル
add_textbox(slide1, "ステップアップナビ",
            Inches(1.5), Inches(1.85), Inches(10.33), Inches(1.0),
            font_size=44, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

# サブタイトル
add_textbox(slide1,
            "障害のある方の『はたらく一歩』を支えるアプリを、一緒に作ってください",
            Inches(1.0), Inches(3.05), Inches(11.33), Inches(0.85),
            font_size=20, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# バッジ
add_badge(slide1, "Readyfor クラウドファンディング プロジェクト",
          Inches(3.7), Inches(6.2), Inches(5.93), Inches(0.55),
          font_size=14)
add_footer(slide1)


# ══════════════════════════════════════════════════════════════
# スライド 2：このプロジェクトで解決したいこと
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK)
body_slide_base(slide2, "このプロジェクトで解決したいこと", "現場で起きているリアルな困りごと")

# リード文カード
add_card(slide2, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.7))
add_textbox(slide2,
            "「自分に何が向いているかわからないまま、就職活動に臨んでいる方がいます。」",
            Inches(0.5), Inches(1.37), Inches(12.3), Inches(0.58),
            font_size=15, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

issues = [
    ("😕  利用者本人の声",
     "「どんな仕事に向いているかわからない」——\n就職先を選ぶ根拠がなく、不安なまま就活に臨んでいます。"),
    ("🤔  支援者（家族・支援員）の声",
     "「何をどう教えればいいか、手探りです」——\n個別に合わせた支援の方法がわからず、疲弊しています。"),
    ("📝  現場の課題",
     "紙とExcelで管理される支援記録——\n情報が共有されず、引き継ぎのたびに記録が途切れます。"),
    ("💡  可能性",
     "テクノロジーを使えば、もっと一人ひとりに合った支援ができるはず。\nそれをカタチにしたのが「ステップアップナビ」です。"),
]

for i, (head, body) in enumerate(issues):
    col = i % 2
    row = i // 2
    lft = Inches(0.35) + col * Inches(6.5)
    top = Inches(2.15) + row * Inches(2.2)
    add_card(slide2, lft, top, Inches(6.1), Inches(2.0))
    add_rect(slide2, lft, top, Inches(6.1), Inches(0.45), fill_color=EMERALD_DARK)
    add_textbox(slide2, head,
                lft + Inches(0.15), top + Inches(0.05), Inches(5.8), Inches(0.38),
                font_size=14, bold=True, color=EMERALD)
    add_multiline_textbox(slide2, body.split("\n"),
                          lft + Inches(0.15), top + Inches(0.55), Inches(5.8), Inches(1.3),
                          font_size=13, color=TEXT_GRAY, line_spacing=Pt(22))


# ══════════════════════════════════════════════════════════════
# スライド 3：解決策 — ステップアップナビとは
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK)
body_slide_base(slide3, "解決策 — ステップアップナビとは", "「はたらく一歩」を一気通貫でサポートするWebアプリ")

# 中央バナー
add_card(slide3, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.8))
add_textbox(slide3,
            "「特性診断 → 成長ロードマップ → 日々の記録 → 支援者との連携」を一気通貫でサポート",
            Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.65),
            font_size=16, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

features = [
    ("🔍", "特性診断",
     "13問のやさしい質問で\n自分の強みと課題を分析"),
    ("🗺️", "成長ロードマップ",
     "仕事タイプ別の3ステップで\n具体的な行動指針を提示"),
    ("🤖", "AIアバターコーチ",
     "自分に似たアバターと会話\n24時間個別アドバイス"),
    ("📊", "支援者連携",
     "家族・支援員が同じ画面で\n情報をリアルタイム共有"),
]

for i, (icon, name, desc) in enumerate(features):
    lft = Inches(0.35) + i * Inches(3.25)
    top = Inches(2.3)
    w = Inches(3.0)
    add_card(slide3, lft, top, w, Inches(2.75))
    add_rect(slide3, lft, top, w, Inches(0.55), fill_color=EMERALD_DARK)
    add_textbox(slide3, f"{icon}  {name}",
                lft + Inches(0.1), top + Inches(0.06), w - Inches(0.2), Inches(0.46),
                font_size=14, bold=True, color=EMERALD)
    add_multiline_textbox(slide3, desc.split("\n"),
                          lft + Inches(0.15), top + Inches(0.7), w - Inches(0.3), Inches(1.8),
                          font_size=13, color=TEXT_GRAY, line_spacing=Pt(22))

# 特長バッジ（下段）
badges = [
    "🌐  スマホ・PCから使える",
    "ひ  難しい言葉を使わない「ひらがなモード」搭載",
    "👤  自分に似たアバターと会話しながら就労を目指せる",
]
add_card(slide3, Inches(0.35), Inches(5.25), Inches(12.6), Inches(1.95))
add_textbox(slide3, "このアプリの3つの特長：",
            Inches(0.5), Inches(5.33), Inches(4), Inches(0.4),
            font_size=13, bold=True, color=EMERALD)
add_multiline_textbox(slide3, badges,
                      Inches(0.5), Inches(5.78), Inches(12.2), Inches(1.3),
                      font_size=13, color=TEXT_GRAY, line_spacing=Pt(24))


# ══════════════════════════════════════════════════════════════
# スライド 4：現在すでに動いているもの（β版）
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK)
body_slide_base(slide4, "現在すでに動いているもの（β版）", "1人のエンジニアが自己資金で開発してきました")

# β版URL強調
add_card(slide4, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.8))
add_rect(slide4, Inches(0.35), Inches(1.3), Inches(0.08), Inches(0.8), fill_color=EMERALD)
add_textbox(slide4, "β版URL：  https://web-production-c0abe.up.railway.app/",
            Inches(0.55), Inches(1.42), Inches(12.3), Inches(0.58),
            font_size=16, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

# バッジ + 見出し
add_badge(slide4, "β版稼働中", Inches(0.35), Inches(2.25), Inches(1.8), Inches(0.4),
          font_size=12)
add_textbox(slide4, "✅  現在実装済みの機能（9機能）",
            Inches(2.3), Inches(2.28), Inches(5), Inches(0.38),
            font_size=14, bold=True, color=EMERALD)

implemented = [
    "✅  特性診断（6特性・13問）",
    "✅  成長ロードマップ（3ステップ）",
    "✅  AIアバターチャット（GPT-4o-mini）",
    "✅  支援者ダッシュボード",
    "✅  ひらがなモード（全ページ対応）",
    "✅  進捗カレンダー・記録機能",
    "✅  日々のタスク管理",
    "✅  感情ログ記録",
    "✅  Railway本番環境稼働中",
]

add_card(slide4, Inches(0.35), Inches(2.8), Inches(6.0), Inches(2.95))
add_multiline_textbox(slide4, implemented[:5],
                      Inches(0.5), Inches(2.95), Inches(5.7), Inches(2.65),
                      font_size=14, color=TEXT_GRAY, line_spacing=Pt(26))
add_card(slide4, Inches(6.65), Inches(2.8), Inches(6.0), Inches(2.95))
add_multiline_textbox(slide4, implemented[5:],
                      Inches(6.8), Inches(2.95), Inches(5.7), Inches(2.65),
                      font_size=14, color=TEXT_GRAY, line_spacing=Pt(26))

# 透明性メッセージ
add_card(slide4, Inches(0.35), Inches(5.95), Inches(12.6), Inches(0.8))
add_rect(slide4, Inches(0.35), Inches(5.95), Inches(0.06), Inches(0.8), fill_color=EMERALD)
add_textbox(slide4,
            "「1人のフリーランス開発者が自己資金で開発してきました。クラウドファンディングで次のステップへ進みたいと思っています。」",
            Inches(0.55), Inches(6.05), Inches(12.3), Inches(0.6),
            font_size=12, bold=True, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 5：開発チーム紹介
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK)
body_slide_base(slide5, "開発チーム紹介", "2人で、現場とテクノロジーをつなぎます")

# 開発者カード（左）
add_card(slide5, Inches(0.35), Inches(1.25), Inches(6.0), Inches(5.7))
add_rect(slide5, Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.6), fill_color=EMERALD_DARK)
add_textbox(slide5, "💻  開発者（松永 直人）",
            Inches(0.5), Inches(1.32), Inches(5.7), Inches(0.48),
            font_size=16, bold=True, color=EMERALD)

dev_items = [
    "フリーランスエンジニア（法人設立検討中）",
    "",
    "担当：Python / Django / AI実装",
    "       OpenAI API × Railway インフラ",
    "",
    "「テクノロジーで社会課題を解決したい」",
    "という想いでゼロから開発を始めました。",
    "",
    "自己資金でβ版を完成させ、現在も継続開発中。",
]
add_multiline_textbox(slide5, dev_items,
                      Inches(0.5), Inches(2.05), Inches(5.7), Inches(4.5),
                      font_size=14, color=TEXT_GRAY, line_spacing=Pt(22))

# パートナーカード（右）
add_card(slide5, Inches(7.0), Inches(1.25), Inches(6.0), Inches(5.7))
add_rect(slide5, Inches(7.0), Inches(1.25), Inches(6.0), Inches(0.6), fill_color=EMERALD_DARK)
add_textbox(slide5, "🏥  パートナー（放課後デイサービス社長）",
            Inches(7.15), Inches(1.32), Inches(5.7), Inches(0.48),
            font_size=15, bold=True, color=EMERALD)

partner_items = [
    "放課後デイサービス運営・社長",
    "",
    "担当：現場の専門家として協力",
    "       ユーザーの声を開発に反映",
    "",
    "実際の支援現場で感じた課題を、",
    "リアルタイムで開発に伝えてくれます。",
    "",
    "「現場で使えるものを一緒に作りたい」",
]
add_multiline_textbox(slide5, partner_items,
                      Inches(7.15), Inches(2.05), Inches(5.7), Inches(4.5),
                      font_size=14, color=TEXT_GRAY, line_spacing=Pt(22))


# ══════════════════════════════════════════════════════════════
# スライド 6：支援金の使い道
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK)
body_slide_base(slide6, "支援金の使い道", "目標金額：300万円 — 1円の無駄なく、現場に届けます")

# 目標金額バッジ
add_badge(slide6, "目標 300万円", Inches(0.35), Inches(1.3), Inches(2.1), Inches(0.5),
          font_size=14)
add_textbox(slide6, "すべて開発・導入・運営に充てます",
            Inches(2.6), Inches(1.35), Inches(10), Inches(0.45),
            font_size=15, color=TEXT_GRAY)

uses = [
    ("AI機能強化（RAG実装）",        "120万円", "40%",
     "個別に最適化されたAIアドバイス機能を実装します"),
    ("スマホ最適化・UI改善",          "60万円",  "20%",
     "スマートフォンでも快適に使えるよう全面改善します"),
    ("サーバー・運営費（1年分）",      "45万円",  "15%",
     "安定したサービス提供のための基盤費用です"),
    ("放課後デイ・事業所テスト導入",   "45万円",  "15%",
     "現場3〜5施設での実証テストと改善サイクルを回します"),
    ("広報・動画制作",                "30万円",  "10%",
     "活動報告・支援者への感謝動画などを制作します"),
]

for i, (name, amount, pct, desc) in enumerate(uses):
    y = Inches(2.0) + i * Inches(0.88)
    add_card(slide6, Inches(0.35), y, Inches(12.6), Inches(0.82))
    add_rect(slide6, Inches(0.35), y, Inches(0.06), Inches(0.82), fill_color=EMERALD)
    add_textbox(slide6, pct,
                Inches(0.52), y + Inches(0.18), Inches(0.8), Inches(0.45),
                font_size=13, bold=True, color=EMERALD)
    add_textbox(slide6, name,
                Inches(1.5), y + Inches(0.08), Inches(4.5), Inches(0.38),
                font_size=14, bold=True, color=TEXT_WHITE)
    add_textbox(slide6, desc,
                Inches(1.5), y + Inches(0.46), Inches(7.5), Inches(0.32),
                font_size=11, color=TEXT_MUTED)
    add_textbox(slide6, amount,
                Inches(10.5), y + Inches(0.15), Inches(2.2), Inches(0.5),
                font_size=20, bold=True, color=EMERALD, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# スライド 7：支援金でできること（マイルストーン）
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK)
body_slide_base(slide7, "支援金でできること（マイルストーン）", "達成額に応じて、できることが広がります")

milestone_header_colors = [EMERALD_DARK, EMERALD_MID, EMERALD]

milestones = [
    ("100万円達成時", [
        "✅  AI個別アドバイス機能（RAG）を実装",
        "✅  支援者評価×多角的分析を追加",
        "✅  基本的なスマホ対応を完了",
    ]),
    ("200万円達成時", [
        "✅  スマートフォン最適化を完了",
        "✅  放課後デイサービス3施設での実証テスト開始",
        "✅  現場フィードバックによる改善",
    ]),
    ("300万円達成時（目標）", [
        "✅  就労支援事業所10施設に無料導入",
        "✅  先輩事例データベース構築開始",
        "✅  1年間の継続運営を保証",
        "✅  開発報告動画・活動報告を継続",
    ]),
]

for i, (title, items) in enumerate(milestones):
    lft = Inches(0.35) + i * Inches(4.35)
    top = Inches(1.25)
    w = Inches(4.05)
    add_card(slide7, lft, top, w, Inches(5.9))
    add_rect(slide7, lft, top, w, Inches(0.65), fill_color=milestone_header_colors[i])
    add_textbox(slide7, title,
                lft + Inches(0.1), top + Inches(0.1), w - Inches(0.2), Inches(0.5),
                font_size=15, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_multiline_textbox(slide7, items,
                          lft + Inches(0.15), top + Inches(0.85), w - Inches(0.3), Inches(4.8),
                          font_size=13, color=TEXT_GRAY, line_spacing=Pt(26))


# ══════════════════════════════════════════════════════════════
# スライド 8：リターン（支援者へのお礼）
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(BLANK)
body_slide_base(slide8, "リターン（支援者へのお礼）", "ご支援いただいた方へ、感謝を込めて")

add_card(slide8, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.65))
add_textbox(slide8,
            "皆さまのご支援が、このアプリを現場に届ける力になります。心よりお礼申し上げます。",
            Inches(0.5), Inches(1.36), Inches(12.3), Inches(0.52),
            font_size=14, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

headers_r = ["支援額", "リターン内容"]
rows_r = [
    ["1,000円",   "お礼メール＋アプリ開発報告（月1回）"],
    ["3,000円",   "支援者クレジット掲載（アプリ内の「つくった人」ページに記名）"],
    ["10,000円",  "アプリ内『応援者一覧ページ』にお名前を永久掲載＋開発進捗レポートをメールで送付（完成まで）"],
    ["30,000円",  "アプリ内『スポンサー』としてお名前・団体名を目立つ形で永久掲載＋現場レポートPDF（年2回）"],
    ["100,000円", "開発へのフィードバック参加権（開発の意思決定に参加）＋事業所向け無料利用（永年）"],
]
add_dark_table(slide8, headers_r, rows_r,
               Inches(0.35), Inches(2.1), Inches(12.6), Inches(4.7),
               font_size=13)

add_card(slide8, Inches(0.35), Inches(6.9), Inches(12.6), Inches(0.4))
add_textbox(slide8, "※ リターンの詳細・送付時期はReadyforプロジェクトページをご確認ください",
            Inches(0.5), Inches(6.93), Inches(12.3), Inches(0.35),
            font_size=11, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════════
# スライド 9：なぜ今クラウドファンディングなのか
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(BLANK)
body_slide_base(slide9, "なぜ今クラウドファンディングなのか", "投資家ではなく、関心を持つ人たちと一緒に作りたい")

# 主メッセージカード
add_card(slide9, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.95))
add_rect(slide9, Inches(0.35), Inches(1.3), Inches(0.08), Inches(0.95), fill_color=EMERALD)
add_textbox(slide9,
            "「まず現場で使われ、現場の声で育てるアプリを作りたい」",
            Inches(0.55), Inches(1.42), Inches(12.3), Inches(0.73),
            font_size=20, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

reasons = [
    ("💬  現場の声で育てたい",
     "投資家の論理ではなく、実際に障害福祉に関心を持つ人たちの声で\nアプリを育てていきたいと考えています。"),
    ("🔨  小さくても動くものがある",
     "β版はすでに稼働中。積み重ねてきた実績があります。\n「できる」と言うだけでなく、「すでに動いている」ことが出発点です。"),
    ("🤝  現場への橋渡しができている",
     "放課後デイサービスとの連携で、開発したものをすぐに現場に届け\nフィードバックを得られる体制が整っています。"),
    ("🌱  一緒に育ててほしい",
     "1人のエンジニアが作ったアプリですが、支援してくださる方の名前を\nアプリの中に刻んで、一緒に育てていきたいと思っています。"),
]

for i, (head, body) in enumerate(reasons):
    col = i % 2
    row = i // 2
    lft = Inches(0.35) + col * Inches(6.5)
    top = Inches(2.5) + row * Inches(2.0)
    add_card(slide9, lft, top, Inches(6.1), Inches(1.85))
    add_rect(slide9, lft, top, Inches(6.1), Inches(0.45), fill_color=EMERALD_DARK)
    add_textbox(slide9, head,
                lft + Inches(0.15), top + Inches(0.05), Inches(5.8), Inches(0.38),
                font_size=13, bold=True, color=EMERALD)
    add_multiline_textbox(slide9, body.split("\n"),
                          lft + Inches(0.15), top + Inches(0.55), Inches(5.8), Inches(1.15),
                          font_size=12, color=TEXT_GRAY, line_spacing=Pt(20))


# ══════════════════════════════════════════════════════════════
# スライド 10：今後のロードマップ
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(BLANK)
body_slide_base(slide10, "今後のロードマップ", "一歩ずつ、着実に現場に届けていきます")

phase_header_colors = [EMERALD_DARK, EMERALD_MID, RGBColor(0x10, 0xA0, 0x70), EMERALD]

phases = [
    ("フェーズ1\n（現在）", [
        "● β版稼働中",
        "● 9機能実装済み",
        "● Railway本番環境",
        "● 放課後デイ連携開始",
    ], "自己資金で開発"),
    ("フェーズ2\n（〜6ヶ月後）", [
        "◎ AI機能強化（RAG）",
        "◎ スマホ最適化完了",
        "◎ 事業所導入テスト",
        "◎ フィードバック収集",
    ], "CF資金で実現"),
    ("フェーズ3\n（〜1年後）", [
        "△ マッチング機能追加",
        "△ 先輩事例データベース",
        "△ 10施設への本格導入",
        "△ 活動報告・開示",
    ], "継続開発フェーズ"),
    ("フェーズ4\n（将来）", [
        "◇ 企業連携・有料化",
        "◇ 全国の事業所へ展開",
        "◇ 障害×就労のDX推進",
        "◇ 社会インフラとして",
    ], "社会実装フェーズ"),
]

for i, (phase, items, tag) in enumerate(phases):
    lft = Inches(0.35) + i * Inches(3.25)
    top = Inches(1.3)
    w = Inches(3.0)
    add_card(slide10, lft, top, w, Inches(5.8))
    add_rect(slide10, lft, top, w, Inches(0.65), fill_color=phase_header_colors[i])
    add_textbox(slide10, phase,
                lft + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.58),
                font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_multiline_textbox(slide10, items,
                          lft + Inches(0.1), top + Inches(0.8), w - Inches(0.2), Inches(3.3),
                          font_size=13, color=TEXT_GRAY, line_spacing=Pt(24))
    add_rect(slide10, lft, top + Inches(4.85), w, Inches(0.5),
             fill_color=phase_header_colors[i])
    add_textbox(slide10, tag,
                lft + Inches(0.05), top + Inches(4.87), w - Inches(0.1), Inches(0.45),
                font_size=11, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# フェーズ間の矢印
for i in range(3):
    lft = Inches(3.2) + i * Inches(3.25)
    add_textbox(slide10, "→",
                lft + Inches(0.02), Inches(3.8), Inches(0.6), Inches(0.6),
                font_size=28, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 11：このアプリが目指す世界
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(BLANK)
add_rect(slide11, 0, 0, W, H, fill_color=BG_DARK)
add_deco_circle(slide11)

# タイトルエリア
add_rect(slide11, Inches(0.3), Inches(0.35), Inches(0.08), Inches(0.5), fill_color=EMERALD)
add_textbox(slide11, "このアプリが目指す世界",
            Inches(0.5), Inches(0.35), Inches(9), Inches(0.5),
            font_size=20, color=TEXT_MUTED)

add_textbox(slide11,
            "障害のある方が、\n自分らしくはたらける社会へ",
            Inches(0.5), Inches(0.95), Inches(9.2), Inches(1.7),
            font_size=36, bold=True, color=TEXT_WHITE)

visions = [
    ("💪", "利用者のために",
     "「障害のある方が自分らしく働ける社会」\n—— 一人ひとりの個性を活かした就労を支えます"),
    ("🤲", "支援者のために",
     "「支援者が迷わず支援できる環境」\n—— 情報共有と記録で、支援者の負担を減らします"),
    ("🌉", "社会のために",
     "「企業と障害者が自然につながる仕組み」\n—— テクノロジーで、社会全体に橋を架けます"),
]

for i, (icon, label, desc) in enumerate(visions):
    lft = Inches(0.5) + i * Inches(4.25)
    top = Inches(3.0)
    add_card(slide11, lft, top, Inches(4.0), Inches(2.9))
    add_rect(slide11, lft, top, Inches(4.0), Inches(0.55), fill_color=EMERALD_DARK)
    add_textbox(slide11, f"{icon}  {label}",
                lft + Inches(0.15), top + Inches(0.08), Inches(3.7), Inches(0.42),
                font_size=15, bold=True, color=EMERALD)
    add_multiline_textbox(slide11, desc.split("\n"),
                          lft + Inches(0.15), top + Inches(0.7), Inches(3.7), Inches(1.9),
                          font_size=13, color=TEXT_GRAY, line_spacing=Pt(22))

add_card(slide11, Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.7))
add_textbox(slide11,
            "一人ひとりの「はたらく一歩」を、テクノロジーで支える",
            Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.55),
            font_size=20, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)
add_footer(slide11)


# ══════════════════════════════════════════════════════════════
# スライド 12：よくある質問（FAQ）
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(BLANK)
body_slide_base(slide12, "よくある質問（FAQ）", "ご不安な点にお答えします")

faqs = [
    ("Q. 開発者1人で大丈夫ですか？",
     "β版はすでに動いています。資金が集まることで開発スピードを上げられます。\n現場パートナーとの連携体制もあり、1人でもリスクを最小化しています。"),
    ("Q. 個人情報は安全ですか？",
     "Railwayの暗号化通信（HTTPS）・パスワードのハッシュ化・\nセキュアなセッション管理で保護しています。"),
    ("Q. 利用者が使いこなせますか？",
     "ひらがなモード・大きなボタン・シンプルUIで設計しています。\n現場の声を反映しながら、より使いやすく改善し続けます。"),
    ("Q. 事業化の見通しはありますか？",
     "放課後デイサービスとの連携でフィードバックを得ながら継続開発します。\nクラウドファンディング後も活動報告を続け、透明性を保ちます。"),
]

for i, (q, a) in enumerate(faqs):
    col = i % 2
    row = i // 2
    lft = Inches(0.35) + col * Inches(6.5)
    top = Inches(1.3) + row * Inches(2.7)
    add_card(slide12, lft, top, Inches(6.1), Inches(2.5))
    add_rect(slide12, lft, top, Inches(6.1), Inches(0.55), fill_color=EMERALD_DARK)
    add_textbox(slide12, q,
                lft + Inches(0.15), top + Inches(0.08), Inches(5.8), Inches(0.42),
                font_size=14, bold=True, color=EMERALD)
    add_multiline_textbox(slide12, a.split("\n"),
                          lft + Inches(0.15), top + Inches(0.7), Inches(5.8), Inches(1.6),
                          font_size=12, color=TEXT_GRAY, line_spacing=Pt(22))

add_card(slide12, Inches(0.35), Inches(6.75), Inches(12.6), Inches(0.5))
add_textbox(slide12,
            "その他ご質問は support@stepupnavi.example.jp までお気軽にどうぞ",
            Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
            font_size=12, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 13：メッセージ（締めくくり）
# ══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(BLANK)
add_rect(slide13, 0, 0, W, H, fill_color=BG_DARK)
add_deco_circle(slide13)

add_rect(slide13, Inches(0.3), Inches(0.3), Inches(0.08), Inches(0.5), fill_color=EMERALD)
add_textbox(slide13, "開発者からのメッセージ",
            Inches(0.5), Inches(0.3), Inches(10), Inches(0.5),
            font_size=18, color=TEXT_MUTED)

# 引用符デコレーション
add_textbox(slide13, "❝",
            Inches(0.5), Inches(1.0), Inches(1.2), Inches(1.0),
            font_size=72, color=EMERALD_DARK, align=PP_ALIGN.LEFT)

# メッセージ本文
message_lines = [
    "一人のエンジニアが、テクノロジーで誰かの",
    "『はたらく一歩』を支えたいと思って作り始めました。",
    "",
    "まだ小さなアプリですが、現場の声とともに育てていきたいと思っています。",
    "",
    "一緒に作ってください。",
]
add_multiline_textbox(slide13, message_lines,
                      Inches(1.3), Inches(1.3), Inches(9.5), Inches(3.5),
                      font_size=22, color=TEXT_WHITE, line_spacing=Pt(34))

add_textbox(slide13, "❞",
            Inches(10.5), Inches(3.8), Inches(1.2), Inches(1.0),
            font_size=72, color=EMERALD_DARK, align=PP_ALIGN.RIGHT)

# 署名
add_rect(slide13, Inches(0.5), Inches(5.0), Inches(4.5), Inches(0.04), fill_color=EMERALD)
add_textbox(slide13, "松永 直人（フリーランスエンジニア）",
            Inches(0.5), Inches(5.1), Inches(6), Inches(0.5),
            font_size=16, bold=True, color=EMERALD)
add_textbox(slide13, "ステップアップナビ  開発者",
            Inches(0.5), Inches(5.6), Inches(6), Inches(0.4),
            font_size=13, color=TEXT_MUTED)

# URL
add_card(slide13, Inches(0.35), Inches(6.3), Inches(12.6), Inches(0.9))
add_textbox(slide13,
            "β版を試してみてください：  https://web-production-c0abe.up.railway.app/",
            Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.65),
            font_size=16, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)
add_footer(slide13)


# ══════════════════════════════════════════════════════════════
# スライド 14：プロジェクト概要・連絡先
# ══════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(BLANK)
add_rect(slide14, 0, 0, W, H, fill_color=BG_DARK)
add_deco_circle(slide14)

add_rect(slide14, Inches(0.3), Inches(0.35), Inches(0.08), Inches(0.5), fill_color=EMERALD)
add_textbox(slide14, "プロジェクト概要・連絡先",
            Inches(0.5), Inches(0.35), Inches(9), Inches(0.5),
            font_size=18, color=TEXT_MUTED)
add_textbox(slide14, "ステップアップナビ",
            Inches(0.5), Inches(0.95), Inches(9), Inches(0.85),
            font_size=40, bold=True, color=TEXT_WHITE)
add_textbox(slide14, "障害のある方の「はたらく一歩」を支えるWebアプリ",
            Inches(0.5), Inches(1.85), Inches(9), Inches(0.5),
            font_size=16, color=EMERALD)

info = [
    ("アプリURL",        "https://web-production-c0abe.up.railway.app/"),
    ("開発者連絡先",     "contact@stepupnavi.example.jp"),
    ("Readyforページ",   "（公開準備中）"),
    ("X（旧Twitter）",   "（TBD）"),
    ("目標金額",         "300万円"),
    ("プロジェクト期間",  "2026年 夏〜秋（予定）"),
]

for i, (k, v) in enumerate(info):
    col = i % 2
    row = i // 2
    lft = Inches(0.5) + col * Inches(6.0)
    top = Inches(2.6) + row * Inches(0.7)
    add_card(slide14, lft, top, Inches(5.8), Inches(0.6))
    add_textbox(slide14, f"{k}：",
                lft + Inches(0.12), top + Inches(0.1), Inches(1.8), Inches(0.42),
                font_size=12, bold=True, color=EMERALD)
    add_textbox(slide14, v,
                lft + Inches(2.1), top + Inches(0.1), Inches(3.5), Inches(0.42),
                font_size=12, color=TEXT_GRAY)

# 下部メッセージ
add_card(slide14, Inches(0.35), Inches(5.5), Inches(12.6), Inches(1.65))
add_textbox(slide14,
            "一緒に作ってください。「はたらく一歩」を、届けましょう。",
            Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.65),
            font_size=22, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)
add_badge(slide14, "Readyfor クラウドファンディング プロジェクト  2026",
          Inches(3.5), Inches(6.35), Inches(6.33), Inches(0.5),
          font_size=13)
add_footer(slide14)


# ─── 保存 ────────────────────────────────────────────────────
OUTPUT = "/Users/matsunaganaoto/Desktop/projects/App/ステップアップナビ_CF資料.pptx"
prs.save(OUTPUT)

slide_count = len(prs.slides)
print(f"✅ クラウドファンディング資料 作成完了（ダークネイビー × エメラルドグリーン リデザイン版）")
print(f"   ファイル：{OUTPUT}")
print(f"   スライド数：{slide_count}枚")
