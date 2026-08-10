#!/usr/bin/env python3
"""東濃信用金庫 提案資料 PPTX生成スクリプト（就労支援×放課後デイ 2軸特化版）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────
# カラー定義
# ──────────────────────────────────────────────
C_NAVY   = RGBColor(0x0f, 0x17, 0x2a)   # 背景
C_GREEN  = RGBColor(0x10, 0xb9, 0x81)   # アクセント1（エメラルドグリーン）
C_BLUE   = RGBColor(0x3b, 0x82, 0xf6)   # アクセント2（ブルー）
C_AMBER  = RGBColor(0xf5, 0x9e, 0x0b)   # 放課後デイサブアクセント
C_WHITE  = RGBColor(0xf8, 0xfa, 0xfc)   # テキスト
C_GRAY   = RGBColor(0x94, 0xa3, 0xb8)   # サブテキスト
C_CARD   = RGBColor(0x1e, 0x29, 0x3b)   # カード背景
C_RED    = RGBColor(0xef, 0x44, 0x44)   # 就労現場アクセント
C_ORANGE = RGBColor(0xea, 0x58, 0x0c)   # 放課後デイ現場アクセント
C_DARK   = RGBColor(0x07, 0x0c, 0x17)   # 深いネイビー

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=C_NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=Pt(13), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def badge(slide, left, top, text, bg_color=C_GREEN, width=Cm(11)):
    h = Cm(0.72)
    rect = add_rect(slide, left, top, width, h, fill_color=bg_color)
    tf = rect.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = C_WHITE


def top_bar(slide, color=C_GREEN):
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.5), fill_color=color)


def bottom_bar(slide):
    add_rect(slide, Cm(0), Cm(18.4), SLIDE_W, Cm(0.65), fill_color=C_DARK)


def accent_line(slide, top, color=C_GREEN, left=Cm(2), width=Cm(6)):
    add_rect(slide, left, top, width, Cm(0.1), fill_color=color)


def slide_header(slide, title, color=C_WHITE, accent_color=C_GREEN,
                 badge_text=None, badge_color=None, title_size=Pt(26)):
    top_bar(slide, color=accent_color)
    if badge_text:
        bc = badge_color if badge_color else accent_color
        badge(slide, Cm(2), Cm(0.8), badge_text, bg_color=bc)
        t_top = Cm(1.9)
    else:
        t_top = Cm(1.0)
    add_textbox(slide, Cm(2), t_top, Cm(29.5), Cm(1.4),
                title, font_size=title_size, bold=True, color=color)
    accent_line(slide, t_top + Cm(1.5), color=accent_color)
    bottom_bar(slide)


# ──────────────────────────────────────────────
# スライド1：表紙
# ──────────────────────────────────────────────
def slide1(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    add_rect(sl, Cm(0), Cm(0), SLIDE_W, Cm(0.6), fill_color=C_GREEN)

    badge(sl, Cm(2), Cm(1.3), "東濃信用金庫様 ご提案資料", bg_color=C_GREEN, width=Cm(12))

    add_textbox(sl, Cm(2), Cm(2.8), Cm(28), Cm(3.0),
                "ステップアップナビ",
                font_size=Pt(58), bold=True, color=C_WHITE)

    accent_line(sl, Cm(6.5), color=C_GREEN, width=Cm(8))

    add_textbox(sl, Cm(2), Cm(7.0), Cm(26), Cm(1.8),
                "障がいのある人が、自分らしく\u300c働き・育つ\u300d社会へ",
                font_size=Pt(20), bold=False, color=C_GRAY)

    add_textbox(sl, Cm(2), Cm(9.2), Cm(26), Cm(1.2),
                "特定非営利活動法人 思いやりの糸 / HIローズ",
                font_size=Pt(13), bold=False, color=C_GRAY)

    for i in range(6):
        x = Cm(23.5) + i * Cm(1.7)
        add_rect(sl, x, Cm(10), Cm(1.3), Cm(8.5),
                 fill_color=RGBColor(0x1e, 0x29, 0x3b))

    add_rect(sl, Cm(0), Cm(18.4), SLIDE_W, Cm(0.65), fill_color=C_DARK)

    add_notes(sl,
        "本日はお時間をいただきありがとうございます。\n"
        "私たちは「ステップアップナビ」というAIを活用した福祉支援アプリを開発しております。\n"
        "就労支援と放課後デイサービスという2つの現場に特化したアプリで、\n"
        "多治見から障がいのある人が自分らしく働き、育てる社会を目指しています。\n"
        "本日は東濃信用金庫様にご支援のご相談にまいりました。どうぞよろしくお願いいたします。"
    )


# ──────────────────────────────────────────────
# スライド2：社会課題（問題提起）
# ──────────────────────────────────────────────
def slide2(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_RED)

    add_textbox(sl, Cm(2), Cm(0.9), Cm(29.5), Cm(1.4),
                "2つの現場が抱える、解決されない問題",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(2.5), color=C_RED)

    C_RED_CARD   = RGBColor(0x2d, 0x10, 0x10)
    C_ORANGE_CARD = RGBColor(0x2a, 0x18, 0x08)

    # ── 左カード：就労現場（赤系） ──
    lx, ly, lw, lh = Cm(1.5), Cm(3.1), Cm(15.2), Cm(14.8)
    add_rect(sl, lx, ly, lw, lh, fill_color=C_RED_CARD)
    add_rect(sl, lx, ly, lw, Cm(0.4), fill_color=C_RED)
    add_textbox(sl, lx + Cm(0.5), ly + Cm(0.7), lw - Cm(1.0), Cm(0.9),
                "【就労現場】", font_size=Pt(18), bold=True, color=C_RED)

    left_items = [
        "企業が障がい者を雇用しても特性を理解できず対処できない",
        "「仕事ができない」レッテルを貼られ解雇になるケースも",
        "企業と障がい者の間に深刻なミスマッチが生じている",
    ]
    for i, item in enumerate(left_items):
        iy = ly + Cm(2.0) + i * Cm(3.8)
        add_rect(sl, lx + Cm(0.5), iy + Cm(0.35), Cm(0.3), Cm(0.3), fill_color=C_RED)
        add_textbox(sl, lx + Cm(1.0), iy, lw - Cm(1.3), Cm(3.2),
                    item, font_size=Pt(13.5), color=C_WHITE, wrap=True)

    # 中央矢印
    add_textbox(sl, Cm(16.5), Cm(9.2), Cm(1.5), Cm(2.0),
                "VS", font_size=Pt(28), bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── 右カード：放課後デイ現場（オレンジ系） ──
    rx, ry, rw, rh = Cm(17.7), Cm(3.1), Cm(14.8), Cm(14.8)
    add_rect(sl, rx, ry, rw, rh, fill_color=C_ORANGE_CARD)
    add_rect(sl, rx, ry, rw, Cm(0.4), fill_color=C_ORANGE)
    add_textbox(sl, rx + Cm(0.5), ry + Cm(0.7), rw - Cm(1.0), Cm(0.9),
                "【放課後デイ現場】", font_size=Pt(18), bold=True, color=C_ORANGE)

    right_items = [
        "児童の特性把握が支援員の経験則に依存している",
        "日々の記録・書類作業に追われ子どもと向き合う時間が不足",
        "保護者への適切なフィードバックができていない",
    ]
    for i, item in enumerate(right_items):
        iy = ry + Cm(2.0) + i * Cm(3.8)
        add_rect(sl, rx + Cm(0.5), iy + Cm(0.35), Cm(0.3), Cm(0.3), fill_color=C_ORANGE)
        add_textbox(sl, rx + Cm(1.0), iy, rw - Cm(1.3), Cm(3.2),
                    item, font_size=Pt(13.5), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "私たちが解決しようとしているのは、2つの現場の深刻な課題です。\n"
        "就労現場では、企業が障がい者を雇用しても特性を理解できず、適切な配慮ができないまま解雇に至るケースが後を絶ちません。\n"
        "放課後デイサービス現場では、支援員が書類作業に追われ、子どもたちと向き合う本来の支援に集中できていない状況があります。\n"
        "これら2つの課題を、テクノロジーの力で同時に解決するのが「ステップアップナビ」です。"
    )


# ──────────────────────────────────────────────
# スライド3：解決アプローチ
# ──────────────────────────────────────────────
def slide3(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_GREEN)

    add_textbox(sl, Cm(2), Cm(0.9), Cm(29.5), Cm(1.4),
                "AIと\u300c寄り添い\u300dで、2つの課題を同時に解決",
                font_size=Pt(25), bold=True, color=C_WHITE)
    accent_line(sl, Cm(2.5), color=C_GREEN)

    C_GREEN_DARK = RGBColor(0x0a, 0x2a, 0x1e)
    C_AMBER_DARK = RGBColor(0x2a, 0x1e, 0x06)

    # 左：就労支援アプリ
    lx, ly, lw, lh = Cm(1.5), Cm(3.2), Cm(15.2), Cm(14.5)
    add_rect(sl, lx, ly, lw, lh, fill_color=C_GREEN_DARK)
    add_rect(sl, lx, ly, lw, Cm(0.45), fill_color=C_GREEN)
    add_textbox(sl, lx + Cm(0.5), ly + Cm(0.75), lw - Cm(1.0), Cm(1.0),
                "就労支援アプリ", font_size=Pt(20), bold=True, color=C_GREEN)

    left_items = [
        ("特性の見える化", "性格・強み・苦手を診断し客観データとして提示"),
        ("企業への配慮提案", "採用企業が「どう配慮すればよいか」を即理解"),
        ("本人の成長サポート", "ロードマップと達成記録で自信とモチベーションを育む"),
    ]
    for i, (head, body) in enumerate(left_items):
        iy = ly + Cm(2.0) + i * Cm(3.8)
        add_rect(sl, lx + Cm(0.4), iy, Cm(0.3), Cm(0.3), fill_color=C_GREEN)
        add_textbox(sl, lx + Cm(1.0), iy - Cm(0.05), lw - Cm(1.3), Cm(0.85),
                    head, font_size=Pt(14), bold=True, color=C_GREEN)
        add_textbox(sl, lx + Cm(1.0), iy + Cm(0.8), lw - Cm(1.3), Cm(2.5),
                    body, font_size=Pt(12.5), color=C_WHITE, wrap=True)

    # 右：放課後デイアプリ
    rx, ry, rw, rh = Cm(17.5), Cm(3.2), Cm(14.8), Cm(14.5)
    add_rect(sl, rx, ry, rw, rh, fill_color=C_AMBER_DARK)
    add_rect(sl, rx, ry, rw, Cm(0.45), fill_color=C_AMBER)
    add_textbox(sl, rx + Cm(0.5), ry + Cm(0.75), rw - Cm(1.0), Cm(1.0),
                "放課後デイアプリ", font_size=Pt(20), bold=True, color=C_AMBER)

    right_items = [
        ("児童の発達記録", "日々の支援データを自動で蓄積・分析"),
        ("支援計画の作成サポート", "支援記録の蓄積をもとに、個別支援計画の作成をサポート・補助し書類作業を大幅軽減"),
        ("保護者連携", "今日のできた！をリアルタイムで保護者と共有"),
    ]
    for i, (head, body) in enumerate(right_items):
        iy = ry + Cm(2.0) + i * Cm(3.8)
        add_rect(sl, rx + Cm(0.4), iy, Cm(0.3), Cm(0.3), fill_color=C_AMBER)
        add_textbox(sl, rx + Cm(1.0), iy - Cm(0.05), rw - Cm(1.3), Cm(0.85),
                    head, font_size=Pt(14), bold=True, color=C_AMBER)
        add_textbox(sl, rx + Cm(1.0), iy + Cm(0.8), rw - Cm(1.3), Cm(2.5),
                    body, font_size=Pt(12.5), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "「ステップアップナビ」は、就労支援と放課後デイサービスという2つのアプリで構成されています。\n"
        "就労支援アプリは、障がいのある大人の方が自分の特性を知り、企業に理解してもらいながら成長できる仕組みを提供します。\n"
        "放課後デイアプリは、支援員が書類作業から解放され、子どもたちと向き合う本当の支援に集中できる環境を作ります。\n"
        "2つのアプリが連携することで、子どもから大人まで一貫した支援の輪が生まれます。"
    )


# ──────────────────────────────────────────────
# スライド3b：AIチャット機能
# ──────────────────────────────────────────────
def slide3b(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_GREEN)

    badge(sl, Cm(2), Cm(0.8), "AIチャット機能", bg_color=C_GREEN, width=Cm(8))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.5),
                "AIが24時間、あなたの\u201c相談相手\u201dになる",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.7), color=C_GREEN)

    cards = [
        (C_GREEN,
         "就労者向け",
         "仕事の悩み・人間関係・体調の変化をチャットで打ち明けると、"
         "自分の特性に合ったアドバイスを返す。支援員がいない時間も孤立させない"),
        (C_AMBER,
         "放課後デイ\n児童・保護者向け",
         "ゲーム感覚でAIキャラクターと会話しながら、自分の気持ちや得意を発見。"
         "保護者は家庭でできる関わり方のヒントをAIから得られる"),
        (C_BLUE,
         "支援員向け",
         "「この子にどう関わればいい？」をAIに聞ける\n\n"
         "・支援の現場では、児童の行動に戸惑い・悩む場面が日々起きる\n"
         "・「こんなとき、どうすればいいか」をAIアドバイザーに相談すると、特性に応じた関わり方のヒントを提案\n"
         "・先輩や専門家に聞きにくい悩みも、AIなら気軽に相談できる\n"
         "・経験の浅い支援員でも、自信を持って支援に臨める\n\n"
         "※支援員の孤立・バーンアウト防止にも貢献"),
    ]

    card_w = Cm(9.5)
    for i, (color, title, body) in enumerate(cards):
        cx = Cm(1.8) + i * Cm(10.4)
        cy = Cm(4.4)
        ch = Cm(12.5)

        add_rect(sl, cx, cy, card_w, ch, fill_color=C_CARD)
        add_rect(sl, cx, cy, card_w, Cm(0.4), fill_color=color)

        add_textbox(sl, cx + Cm(0.5), cy + Cm(0.7), card_w - Cm(1.0), Cm(1.5),
                    title, font_size=Pt(17), bold=True, color=color, wrap=True)

        accent_line(sl, cy + Cm(2.5), color=color, left=cx + Cm(0.5), width=card_w - Cm(1.0))

        add_textbox(sl, cx + Cm(0.5), cy + Cm(3.0), card_w - Cm(1.0), Cm(8.5),
                    body, font_size=Pt(13), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "ステップアップナビの中核機能の一つが、24時間対応のAIチャットです。\n"
        "就労者は、仕事の悩みや体調の変化を気軽にチャットで打ち明けられます。自分の特性に合ったアドバイスが返ってくるため、支援員が不在の時間も孤立しません。\n"
        "放課後デイの児童・保護者には、ゲーム感覚のAIキャラクターとの会話を通じて、子ども自身の気持ちや得意の発見を促します。保護者も家庭での関わり方のヒントをAIから得られます。\n"
        "支援員は、困難なケースについてAIに相談することで、対処法のヒントや精神的な支えを得られます。「人には言いにくいことも、AIなら話せる」という安心感が現場を支えます。"
    )


# ──────────────────────────────────────────────
# スライド4：就労支援アプリ詳細
# ──────────────────────────────────────────────
def slide4(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_GREEN)

    badge(sl, Cm(2), Cm(0.8), "就労支援（大人向け）", bg_color=C_GREEN, width=Cm(9.5))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.5),
                "\u300cできない\u300dを\u300cできる\u300dに変える、成長の地図",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.7), color=C_GREEN)

    items = [
        ("性格・特性診断",
         "「得意なこと・苦手なこと」を可視化\n本人も支援者も企業担当者も同じ地図を共有できる"),
        ("3ステップのロードマップ",
         "仕事に必要なスキルを段階的に提示\n「今日やること」が明確になり、迷わず前進できる"),
        ("達成記録で自信を育む",
         "日々の小さな達成を記録し、モチベーションと自己肯定感を継続的に育む"),
        ("支援者・企業ダッシュボード",
         "本人の状況を支援者・企業担当者がリアルタイムで把握\nミスマッチを事前に防ぐ"),
    ]

    for i, (head, body) in enumerate(items):
        row = i // 2
        col = i % 2
        cx = Cm(2.0) + col * Cm(15.7)
        cy = Cm(4.3) + row * Cm(6.6)
        cw, ch = Cm(14.8), Cm(6.0)

        add_rect(sl, cx, cy, cw, ch, fill_color=C_CARD)
        add_rect(sl, cx, cy, Cm(0.2), ch, fill_color=C_GREEN)

        add_textbox(sl, cx + Cm(0.6), cy + Cm(0.4), cw - Cm(0.9), Cm(0.95),
                    head, font_size=Pt(15.5), bold=True, color=C_GREEN)
        add_textbox(sl, cx + Cm(0.6), cy + Cm(1.45), cw - Cm(0.9), Cm(4.0),
                    body, font_size=Pt(13), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "就労支援アプリの核心は「成長の地図」です。\n"
        "まず性格・特性診断で、本人の得意なこと・苦手なことを客観的に可視化します。\n"
        "次に、仕事に必要なスキルを3段階のロードマップで提示し、今日やることを明確にします。\n"
        "日々の達成記録が積み重なることで、自信とモチベーションが育まれます。\n"
        "さらに支援者・企業担当者向けダッシュボードで、本人の状況をリアルタイム共有し、ミスマッチを防ぎます。"
    )


# ──────────────────────────────────────────────
# スライド5：企業側へのメリット
# ──────────────────────────────────────────────
def slide5(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_BLUE)

    badge(sl, Cm(2), Cm(0.8), "採用企業向け", bg_color=C_BLUE, width=Cm(7))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.5),
                "企業が「採用してよかった」と思える仕組み",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.7), color=C_BLUE)

    items = [
        "本人の特性レポートで「どう配慮すれば良いか」が一目でわかる",
        "成長記録を共有することで、長期雇用につながる信頼関係を構築",
        "法定雇用率の達成と定着率向上を同時に実現",
    ]

    for i, item in enumerate(items):
        y = Cm(4.5) + i * Cm(4.0)
        add_rect(sl, Cm(2), y, Cm(29.5), Cm(3.5), fill_color=C_CARD)
        add_rect(sl, Cm(2), y, Cm(0.22), Cm(3.5), fill_color=C_BLUE)

        num_bg = add_rect(sl, Cm(2.5), y + Cm(0.8), Cm(1.2), Cm(1.2), fill_color=C_BLUE)
        tf_nb = num_bg.text_frame
        p_nb = tf_nb.paragraphs[0]
        p_nb.alignment = PP_ALIGN.CENTER
        rn = p_nb.add_run()
        rn.text = str(i + 1)
        rn.font.size = Pt(16)
        rn.font.bold = True
        rn.font.color.rgb = C_WHITE

        add_textbox(sl, Cm(4.2), y + Cm(0.85), Cm(26.8), Cm(1.7),
                    item, font_size=Pt(14.5), color=C_WHITE, wrap=True)

    add_textbox(sl, Cm(2), Cm(17.0), Cm(29.5), Cm(1.0),
                "障がい者雇用の「採用→定着→成長」を一つのアプリで完結",
                font_size=Pt(14), bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    bottom_bar(sl)

    add_notes(sl,
        "このアプリは障がい者本人だけでなく、採用企業にも大きなメリットをもたらします。\n"
        "特性レポートを見れば、「この人にはどう配慮すれば良いか」が採用前からわかります。\n"
        "成長記録の共有により、企業と障がい者の間に継続的な信頼関係が生まれ、長期雇用につながります。\n"
        "法定雇用率の達成と定着率向上を同時に実現できる、企業にとっても価値あるツールです。"
    )


# ──────────────────────────────────────────────
# スライド6：放課後デイサービスアプリ詳細
# ──────────────────────────────────────────────
def slide6(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_AMBER)

    badge(sl, Cm(2), Cm(0.8), "放課後デイサービス向け", bg_color=C_AMBER, width=Cm(10))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.5),
                "支援員が\u300c本当にやりたい支援\u300dに集中できる環境へ",
                font_size=Pt(25), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.7), color=C_AMBER)

    items = [
        ("音声入力＋AI効率化",
         "日々の支援記録を音声入力＋AIで効率化・補助\n書類作業を大幅削減し、子どもとの時間を最大化"),
        ("個別支援計画の作成サポート",
         "支援記録のデータをもとに、個別支援計画の作成をサポート・補助\n経験則に頼らない客観的な支援が実現"),
        ("保護者アプリと連携",
         "「今日のできた！」をリアルタイムで保護者と共有\n保護者の安心感と事業所への信頼が高まる"),
        ("発達の見える化",
         "発達の変化を時系列グラフで可視化\n支援の成果が「エビデンス」として蓄積される"),
    ]

    for i, (head, body) in enumerate(items):
        row = i // 2
        col = i % 2
        cx = Cm(2.0) + col * Cm(15.7)
        cy = Cm(4.3) + row * Cm(6.6)
        cw, ch = Cm(14.8), Cm(6.0)

        add_rect(sl, cx, cy, cw, ch, fill_color=C_CARD)
        add_rect(sl, cx, cy, Cm(0.2), ch, fill_color=C_AMBER)

        add_textbox(sl, cx + Cm(0.6), cy + Cm(0.4), cw - Cm(0.9), Cm(0.95),
                    head, font_size=Pt(15.5), bold=True, color=C_AMBER)
        add_textbox(sl, cx + Cm(0.6), cy + Cm(1.45), cw - Cm(0.9), Cm(4.0),
                    body, font_size=Pt(13), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "放課後デイサービスアプリは、支援員が本来の仕事に集中できる環境づくりを目指しています。\n"
        "音声入力とAIの組み合わせで、日々の支援記録の入力が効率化・補助されます。書類作業の時間が大幅に減ります。\n"
        "支援記録のデータをもとにAIが個別支援計画の作成をサポートし、科学的な支援が可能になります。\n"
        "保護者アプリとの連携で「今日のできた！」をリアルタイム共有でき、保護者との信頼関係が深まります。\n"
        "発達の変化を時系列グラフで見える化することで、支援の成果が明確になります。"
    )


# ──────────────────────────────────────────────
# スライド7：開発者の思い
# ──────────────────────────────────────────────
def slide7(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_GREEN)

    add_textbox(sl, Cm(2), Cm(0.9), Cm(29.5), Cm(1.4),
                "当事者家族だから作れる、\u300c本当に使えるアプリ\u300d",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(2.5), color=C_GREEN)

    story_items = [
        ("👨‍👧 当事者家族として",
         "開発者は障がいを持つ娘の父親であり、フリーランスエンジニア\n「既存のシステムは現場の声が届いていない」という実感から開発を決意"),
        ("🤝 現場との共同開発",
         "放課後デイサービスの社長とタッグを組み、現場の声を直接反映\n「作りたいもの」ではなく「本当に必要なもの」を追求"),
        ("🚀 フルコミットで開発中",
         "7月〜9月の3ヶ月間、フルコミットで開発中（現在60%完成済み）\n娘のため・地域のために命がけで取り組む"),
    ]

    for i, (head, body) in enumerate(story_items):
        y = Cm(3.2) + i * Cm(4.8)
        add_rect(sl, Cm(2), y, Cm(29.5), Cm(4.3), fill_color=C_CARD)
        add_rect(sl, Cm(2), y, Cm(0.22), Cm(4.3), fill_color=C_GREEN)

        add_textbox(sl, Cm(2.7), y + Cm(0.35), Cm(28.5), Cm(0.95),
                    head, font_size=Pt(15), bold=True, color=C_GREEN)
        add_textbox(sl, Cm(2.7), y + Cm(1.35), Cm(28.5), Cm(2.6),
                    body, font_size=Pt(13), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "このアプリが「本当に使える」理由は、開発者自身が当事者家族だからです。\n"
        "障がいを持つ娘を持つ父として、また放課後デイサービスの現場社長とタッグを組んだことで、\n"
        "既存システムが抱える「現場の声が届いていない」という課題を直接解決しています。\n"
        "現在7月〜9月の3ヶ月間、フルコミットで開発を進めており、すでに60%が完成しています。\n"
        "娘のため、そして地域の皆さんのために、命がけで取り組んでいます。"
    )


# ──────────────────────────────────────────────
# スライド8：機能ロードマップ
# ──────────────────────────────────────────────
def slide8(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_BLUE)

    add_textbox(sl, Cm(2), Cm(0.9), Cm(29.5), Cm(1.4),
                "段階的な開発で、確実に価値を届ける",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(2.5), color=C_BLUE)

    phases = [
        ("Phase 1", "〜3ヶ月", "基盤構築 + 主要機能 + 地域展開", C_GREEN,
         [
             "性格・特性診断とロードマップ機能",
             "日々の記録・成長の可視化",
             "アバターによるAIチャット相談機能",
             "支援記録の音声入力・効率化サポート",
             "就労者向け・放課後デイ向け 基本機能すべてリリース",
             "企業向けダッシュボード（特性レポート・定着サポート）",
             "保護者連携アプリ",
             "地域企業・行政とのタイアップ",
             "他自治体へのライセンス展開",
         ],
         "全基本機能リリース＋\n地域展開スタート"),
        ("Phase 2", "〜1年", "拡張・地域展開", C_BLUE,
         [
             "企業向けダッシュボード\n（特性レポート・定着サポート）",
             "保護者連携アプリ",
             "地域企業・行政とのタイアップ",
             "他自治体へのライセンス展開",
         ],
         "地域展開・\nライセンス全国展開"),
    ]

    phase_w = Cm(14.9)
    card_h = Cm(15.0)

    for i, (phase, period, subtitle, color, feats, summary) in enumerate(phases):
        x = Cm(1.5) + i * (phase_w + Cm(1.0))
        y = Cm(3.2)

        # Phase1（9項目）は小さいフォント・狭いスペーシング、Phase2は通常
        feat_font = Pt(11) if len(feats) > 5 else Pt(13.5)
        feat_spacing = Cm(1.0) if len(feats) > 5 else Cm(1.85)
        feat_h = Cm(0.95) if len(feats) > 5 else Cm(1.75)

        add_rect(sl, x, y, phase_w, card_h, fill_color=C_CARD)
        add_rect(sl, x, y, phase_w, Cm(0.45), fill_color=color)

        add_textbox(sl, x + Cm(0.5), y + Cm(0.7), phase_w - Cm(1.0), Cm(0.95),
                    phase, font_size=Pt(22), bold=True, color=color)
        add_textbox(sl, x + Cm(0.5), y + Cm(1.65), phase_w - Cm(1.0), Cm(0.6),
                    period, font_size=Pt(12.5), bold=False, color=C_GRAY)
        add_textbox(sl, x + Cm(0.5), y + Cm(2.25), phase_w - Cm(1.0), Cm(0.7),
                    subtitle, font_size=Pt(13), bold=True, color=C_WHITE)

        accent_line(sl, y + Cm(3.05), color=color, left=x + Cm(0.5), width=phase_w - Cm(1.0))

        for j, feat in enumerate(feats):
            fy = y + Cm(3.55) + j * feat_spacing
            add_rect(sl, x + Cm(0.6), fy + Cm(0.4), Cm(0.25), Cm(0.25), fill_color=color)
            add_textbox(sl, x + Cm(1.15), fy, phase_w - Cm(1.5), feat_h,
                        feat, font_size=feat_font, color=C_WHITE, wrap=True)

        sum_y = y + Cm(12.6)
        add_rect(sl, x + Cm(0.5), sum_y, phase_w - Cm(1.0), Cm(2.0),
                 fill_color=RGBColor(0x07, 0x0c, 0x17))
        add_textbox(sl, x + Cm(0.6), sum_y + Cm(0.25), phase_w - Cm(1.2), Cm(1.6),
                    summary, font_size=Pt(13), bold=True, color=color,
                    align=PP_ALIGN.CENTER, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "開発は2つのフェーズで段階的に進めます。\n"
        "Phase1（〜3ヶ月）では、性格・特性診断・ロードマップ・日々の記録・アバターAIチャット・音声入力など基本機能に加え、企業向けダッシュボード・保護者連携アプリ・地域企業や行政とのタイアップ・他自治体へのライセンス展開まで全9機能をリリースします。\n"
        "Phase2（〜1年）では、Phase1で構築した基盤をさらに拡張・深化させ、地域全体・全国へ広げていきます。"
    )


# ──────────────────────────────────────────────
# スライド9：東濃信用金庫様へのご提案
# ──────────────────────────────────────────────
def slide9(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_BLUE)

    badge(sl, Cm(2), Cm(0.8), "東濃信用金庫様へのご提案", bg_color=C_BLUE, width=Cm(12))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.5),
                "地域の未来を、一緒に作りませんか",
                font_size=Pt(28), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.7), color=C_BLUE)

    bullets = [
        ("Readyforクラウドファンディングの手数料が無料に",
         "通常8%の手数料が信金パートナーとしてご参画いただくことで0円に"),
        ("300万円目標 → 約24万円のコスト削減",
         "削減分がすべてアプリ開発・普及活動に充当される"),
        ("多治見発・地域貢献の先進事例として信金の名前を全国に発信",
         "福祉DXのパイオニアとして東濃信用金庫様のブランドを広くPR"),
    ]

    for i, (head, sub) in enumerate(bullets):
        y = Cm(4.5) + i * Cm(3.8)
        add_rect(sl, Cm(2), y, Cm(21), Cm(3.3), fill_color=C_CARD)
        add_rect(sl, Cm(2), y, Cm(0.22), Cm(3.3), fill_color=C_BLUE)

        add_textbox(sl, Cm(2.7), y + Cm(0.25), Cm(19.8), Cm(1.0),
                    head, font_size=Pt(14.5), bold=True, color=C_WHITE)
        add_textbox(sl, Cm(2.7), y + Cm(1.3), Cm(19.8), Cm(1.7),
                    sub, font_size=Pt(12.5), color=C_GRAY, wrap=True)

    # 強調カード
    card_x = Cm(24.5)
    add_rect(sl, card_x, Cm(4.5), Cm(8), Cm(11), fill_color=C_BLUE)
    add_textbox(sl, card_x, Cm(5.5), Cm(8), Cm(1.2),
                "削減手数料", font_size=Pt(13), bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, card_x, Cm(7.0), Cm(8), Cm(2.0),
                "¥240,000", font_size=Pt(34), bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, card_x, Cm(9.3), Cm(8), Cm(1.0),
                "がプロジェクトに", font_size=Pt(13),
                color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, card_x, Cm(10.5), Cm(8), Cm(1.0),
                "フル活用されます", font_size=Pt(13),
                color=C_WHITE, align=PP_ALIGN.CENTER)

    bottom_bar(sl)

    add_notes(sl,
        "東濃信用金庫様への具体的なご提案です。\n"
        "Readyforには、信用金庫がパートナーとして参画することでクラウドファンディング手数料が無料になる制度があります。\n"
        "300万円の目標であれば、約24万円の手数料がすべてプロジェクトに活きます。\n"
        "多治見発の地域貢献プロジェクトを支援した先進事例として、東濃信用金庫様のお名前を全国に発信できます。"
    )


# ──────────────────────────────────────────────
# スライド10：300万円の根拠
# ──────────────────────────────────────────────
def slide10(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_GREEN)

    add_textbox(sl, Cm(2), Cm(0.9), Cm(29.5), Cm(1.4),
                "クラウドファンディング 300万円の根拠",
                font_size=Pt(26), bold=True, color=C_WHITE)
    accent_line(sl, Cm(2.5), color=C_GREEN)

    add_textbox(sl, Cm(2), Cm(3.1), Cm(29.5), Cm(0.8),
                "エンジニア単価 約83万円/月 × 3ヶ月 = 250万円 ＋ 運営費50万円 = 合計300万円",
                font_size=Pt(17), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # ━━━ 左：メインテーブル（合計300万円） ━━━
    rows = [
        ("開発人件費（約83万×3ヶ月）",   "250万円", "設計・実装・テスト・リリース・保守"),
        ("AWSインフラ費（3ヶ月分）",      "  3万円", "EC2・RDS・S3・CloudFront"),
        ("広報・PR費",                    " 20万円", "SNS広告、プレスリリース、チラシ制作"),
        ("デザイン・UX制作費",            " 12万円", "UI改善、アイコン、動画制作"),
        ("イベント・説明会開催費",         "  9万円", "福祉事業所向け説明会、デモ展示"),
        ("福祉事業所導入サポート費",       "  6万円", "初期研修・マニュアル・問い合わせ対応"),
    ]

    table_x = Cm(1.5)
    table_w = Cm(18.5)
    row_h   = Cm(1.2)
    row_top = Cm(4.45)

    # ヘッダー行
    add_rect(sl, table_x, row_top - Cm(0.45), table_w, Cm(0.45), fill_color=C_GREEN)
    add_textbox(sl, table_x + Cm(0.4), row_top - Cm(0.4), table_w - Cm(0.5), Cm(0.4),
                "  項目                                金額        内容",
                font_size=Pt(10), bold=True, color=C_WHITE)

    for idx, (label, amount, detail) in enumerate(rows):
        y = row_top + idx * row_h
        bg_color = C_CARD if idx % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
        add_rect(sl, table_x, y, table_w, row_h - Cm(0.05), fill_color=bg_color)
        add_rect(sl, table_x, y, Cm(0.18), row_h - Cm(0.05), fill_color=C_GREEN)

        add_textbox(sl, table_x + Cm(0.4), y + Cm(0.15), Cm(9.2), Cm(0.9),
                    label, font_size=Pt(11), color=C_WHITE)
        add_textbox(sl, table_x + Cm(9.7), y + Cm(0.15), Cm(3.0), Cm(0.9),
                    amount, font_size=Pt(11), bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)
        add_textbox(sl, table_x + Cm(13.0), y + Cm(0.15), Cm(5.0), Cm(0.9),
                    detail, font_size=Pt(9.5), color=C_GRAY)

    # 合計行
    total_y = row_top + len(rows) * row_h + Cm(0.1)
    add_rect(sl, table_x, total_y, table_w, row_h - Cm(0.05), fill_color=C_GREEN)
    add_textbox(sl, table_x + Cm(0.4), total_y + Cm(0.15), Cm(9.2), Cm(0.9),
                "合計", font_size=Pt(12), bold=True, color=C_WHITE)
    add_textbox(sl, table_x + Cm(9.7), total_y + Cm(0.15), Cm(3.0), Cm(0.9),
                "300万円", font_size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

    # フットノート
    fn_y = total_y + row_h + Cm(0.1)
    add_textbox(sl, table_x, fn_y, table_w, Cm(1.2),
                "※AWSは月約1万円〜のスモールスタート。利用者増加に応じてスケールアップ可能。\n"
                "※フリーランスエンジニア市場単価：月60〜120万円が相場。AI×福祉の専門性を持つエンジニアは希少。",
                font_size=Pt(9), bold=False, color=C_GRAY, wrap=True)

    # ━━━ 右：開発人件費250万円の内訳ボックス ━━━
    box_x   = Cm(21.0)
    box_w   = Cm(12.0)
    box_top = Cm(4.0)
    box_h   = Cm(14.0)

    add_rect(sl, box_x, box_top, box_w, box_h, fill_color=C_CARD)
    add_rect(sl, box_x, box_top, box_w, Cm(0.4), fill_color=C_GREEN)

    add_textbox(sl, box_x + Cm(0.3), box_top + Cm(0.55), box_w - Cm(0.5), Cm(0.65),
                "▼ 開発人件費 250万円の内訳",
                font_size=Pt(12), bold=True, color=C_GREEN)

    add_textbox(sl, box_x + Cm(0.3), box_top + Cm(1.25), box_w - Cm(0.5), Cm(0.55),
                "「なぜ月83万円かかるのか？」への回答",
                font_size=Pt(10), bold=False, color=C_GRAY)

    # サブテーブル ヘッダー
    sh_y = box_top + Cm(1.95)
    add_rect(sl, box_x + Cm(0.2), sh_y, box_w - Cm(0.4), Cm(0.42),
             fill_color=RGBColor(0x0a, 0x55, 0x3c))
    add_textbox(sl, box_x + Cm(0.35), sh_y + Cm(0.04), Cm(3.6), Cm(0.36),
                "工程", font_size=Pt(9), bold=True, color=C_WHITE)
    add_textbox(sl, box_x + Cm(4.05), sh_y + Cm(0.04), Cm(1.9), Cm(0.36),
                "期間", font_size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, box_x + Cm(6.1), sh_y + Cm(0.04), Cm(5.5), Cm(0.36),
                "内容", font_size=Pt(9), bold=True, color=C_WHITE)

    # サブテーブル データ行
    sub_rows = [
        ("要件定義・設計",    "2週間", "システム設計、DB設計、UI/UXワイヤーフレーム"),
        ("バックエンド開発",  "4週間", "Django/Python、API実装、AI連携（OpenAI）"),
        ("フロントエンド開発","3週間", "画面実装、レスポンシブ対応、アバター機能"),
        ("テスト・品質保証",  "2週間", "単体テスト、結合テスト、セキュリティ確認"),
        ("リリース・導入支援","1週間", "本番環境構築、マニュアル作成、初期サポート"),
    ]

    sub_row_h = Cm(1.22)
    sr_top = sh_y + Cm(0.42)
    for i, (process, duration, content) in enumerate(sub_rows):
        sy = sr_top + i * sub_row_h
        bg = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
        add_rect(sl, box_x + Cm(0.2), sy, box_w - Cm(0.4), sub_row_h - Cm(0.04), fill_color=bg)

        add_textbox(sl, box_x + Cm(0.35), sy + Cm(0.08), Cm(3.6), Cm(1.1),
                    process, font_size=Pt(9), bold=True, color=C_WHITE, wrap=True)
        add_textbox(sl, box_x + Cm(4.05), sy + Cm(0.08), Cm(1.9), Cm(1.1),
                    duration, font_size=Pt(9), color=C_GREEN, align=PP_ALIGN.CENTER)
        add_textbox(sl, box_x + Cm(6.1), sy + Cm(0.08), Cm(5.5), Cm(1.1),
                    content, font_size=Pt(8.5), color=C_GRAY, wrap=True)

    # 市場単価セクション
    mkt_top = sr_top + len(sub_rows) * sub_row_h + Cm(0.22)
    add_rect(sl, box_x + Cm(0.2), mkt_top, box_w - Cm(0.4), Cm(0.38),
             fill_color=RGBColor(0x07, 0x0c, 0x17))
    add_textbox(sl, box_x + Cm(0.35), mkt_top + Cm(0.04), box_w - Cm(0.6), Cm(0.32),
                "フリーランス市場単価の根拠", font_size=Pt(9), bold=True, color=C_AMBER)

    market_items = [
        "・AI開発経験者：月80〜120万円が相場（大手求人サイト調べ）",
        "・福祉×IT専門性：希少領域のため高単価",
        "・社会保険・経費込みで会社員換算 年収1,000万円相当",
    ]
    for i, item in enumerate(market_items):
        add_textbox(sl, box_x + Cm(0.35),
                    mkt_top + Cm(0.42) + i * Cm(0.72),
                    box_w - Cm(0.6), Cm(0.68),
                    item, font_size=Pt(8.5), color=C_GRAY, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "こちらが資金の使い道です。\n"
        "メインテーブルの通り、開発人件費250万円・AWSインフラ3万円・広報PR20万円・デザイン12万円・イベント9万円・導入サポート6万円の合計300万円です。\n"
        "右側の内訳ボックスが「なぜ月83万円かかるのか」への回答です。\n"
        "要件定義〜リリース・導入支援まで合計12週間（3ヶ月）のフルコミット工程で構成されており、\n"
        "AI×福祉領域の専門エンジニアの市場単価は月80〜120万円が相場です。\n"
        "娘のため、そして地域の皆さんのために開発に専念する計画です。"
    )


# ──────────────────────────────────────────────
# スライド11：スポンサー参画のメリット
# ──────────────────────────────────────────────
def slide11(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    top_bar(sl, color=C_BLUE)

    badge(sl, Cm(2), Cm(0.8), "東濃信用金庫様限定", bg_color=C_BLUE, width=Cm(10))

    add_textbox(sl, Cm(2), Cm(2.0), Cm(29.5), Cm(1.4),
                "スポンサー参画のメリット",
                font_size=Pt(28), bold=True, color=C_WHITE)
    accent_line(sl, Cm(3.6), color=C_BLUE)

    merits = [
        ("ブランド価値",
         "地域福祉DXのパイオニアとして\nメディア・SNSで発信\nアプリ起動画面に信金名掲載"),
        ("地域連携",
         "福祉事業所・保護者・行政との\nネットワーク強化\n地元企業・住民との信頼構築"),
        ("PR効果",
         "アプリ内・クラファンページ・\n広報資料への信金名掲載\n全国の注目案件として拡散"),
        ("社会貢献\n（ESG・CSR）",
         "障がい者就労・児童発達支援\nという社会的インパクト\nESG・CSR評価の向上に貢献"),
    ]

    card_w = Cm(7.3)
    for i, (title, desc) in enumerate(merits):
        x = Cm(1.7) + i * Cm(7.9)
        y = Cm(4.3)
        add_rect(sl, x, y, card_w, Cm(12.5), fill_color=C_CARD)
        add_rect(sl, x, y, card_w, Cm(0.35), fill_color=C_BLUE)

        add_textbox(sl, x + Cm(0.4), y + Cm(0.7), card_w - Cm(0.8), Cm(1.5),
                    title, font_size=Pt(16), bold=True, color=C_BLUE, wrap=True)
        add_textbox(sl, x + Cm(0.4), y + Cm(2.5), card_w - Cm(0.8), Cm(8.5),
                    desc, font_size=Pt(12.5), color=C_WHITE, wrap=True)

    bottom_bar(sl)

    add_notes(sl,
        "スポンサーとしてご参画いただくことで、4つの大きなメリットがあります。\n"
        "ブランド価値：地域福祉DXのパイオニアとして信金の名前を広く発信します。\n"
        "地域連携：福祉事業所・保護者・行政機関との新たなネットワーク構築につながります。\n"
        "PR効果：アプリ内や広報資料への掲載で全国に発信されます。\n"
        "社会貢献：障がい者就労支援・児童発達支援というESG・CSRに直結する取り組みです。"
    )


# ──────────────────────────────────────────────
# スライド12：締めくくり
# ──────────────────────────────────────────────
def slide12(prs):
    sl = add_blank_slide(prs)
    set_bg(sl)
    add_rect(sl, Cm(0), Cm(0), SLIDE_W, Cm(0.6), fill_color=C_GREEN)

    add_textbox(sl, Cm(2), Cm(1.5), Cm(29.5), Cm(2.0),
                "「共犯者」になってくれませんか？",
                font_size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(10), Cm(4.2), Cm(14), Cm(0.12), fill_color=C_GREEN)

    card = add_rect(sl, Cm(2.5), Cm(4.9), Cm(28.5), Cm(8.5), fill_color=C_CARD)
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "障がいのある人も、その家族も、支援する人も——\n\n"
        "みんながもっと笑顔になれる社会を、多治見から作ります。\n\n"
        "まず一度、お話しを聞いていただけますか。"
    )
    run.font.size = Pt(17)
    run.font.color.rgb = C_WHITE

    add_textbox(sl, Cm(2), Cm(14.5), Cm(29.5), Cm(1.0),
                "特定非営利活動法人 思いやりの糸 / HIローズ",
                font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(sl, Cm(2), Cm(15.6), Cm(29.5), Cm(1.0),
                "代表 廣瀬 豊",
                font_size=Pt(14), bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(0), Cm(18.4), SLIDE_W, Cm(0.65), fill_color=C_DARK)

    add_notes(sl,
        "最後に一言お伝えさせてください。\n"
        "これは一人の父が、障がいを持つ娘のために始めたプロジェクトです。\n"
        "しかし、このアプリが完成すれば、多治見だけでなく全国の障がいのある人たちの自立を支えるインフラになり得ます。\n"
        "障がいのある人も、その家族も、支援する人も——みんながもっと笑顔になれる社会を一緒に作りませんか。\n"
        "東濃信用金庫様も、ぜひ「共犯者」になっていただけますよう、お願い申し上げます。\n"
        "本日はありがとうございました。"
    )


# ──────────────────────────────────────────────
# メイン
# ──────────────────────────────────────────────
def main():
    prs = new_prs()

    slide1(prs)   # 1.  表紙
    slide2(prs)   # 2.  社会課題（問題提起）
    slide3(prs)   # 3.  解決アプローチ
    slide3b(prs)  # 4.  AIチャット機能（新規）
    slide4(prs)   # 5.  就労支援アプリ詳細
    slide5(prs)   # 6.  企業側へのメリット
    slide6(prs)   # 7.  放課後デイサービスアプリ詳細
    slide7(prs)   # 8.  開発者の思い
    slide8(prs)   # 9.  機能ロードマップ
    slide9(prs)   # 10. 東濃信用金庫様へのご提案
    slide10(prs)  # 11. 300万円の根拠
    slide11(prs)  # 12. スポンサー参画のメリット
    slide12(prs)  # 13. 締めくくり

    output = "/Users/matsunaganaoto/Desktop/projects/App/東濃信用金庫_提案資料.pptx"
    prs.save(output)
    print(f"✅ 生成完了: {output}")
    print(f"   スライド数: {len(prs.slides)} 枚")


if __name__ == "__main__":
    main()
