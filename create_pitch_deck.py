"""
ステップアップナビ ピッチデック生成スクリプト
投資家・スポンサー向け（全15枚）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree

# ─── カラーパレット ────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x1B, 0x45, 0x32)
MID_GREEN   = RGBColor(0x2D, 0x6A, 0x4F)
LIGHT_GREEN = RGBColor(0x95, 0xD5, 0xB2)
MINT        = RGBColor(0xD8, 0xF3, 0xDC)
ORANGE      = RGBColor(0xF4, 0xA2, 0x61)
ORANGE_DARK = RGBColor(0xE7, 0x6F, 0x51)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xF8)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY    = RGBColor(0x6C, 0x75, 0x7D)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # 完全空白レイアウト


# ─── ヘルパー関数 ─────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=16, bold=False, color=DARK_TEXT,
                          align=PP_ALIGN.LEFT, font_name="Meiryo",
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def body_slide_base(slide, title_text, subtitle_text=None):
    """本文スライドの共通レイアウト: OFF_WHITE背景 + 左MID_GREEN縦ライン + タイトル"""
    # 背景
    add_rect(slide, 0, 0, W, H, fill_color=OFF_WHITE)
    # 左縦ライン
    add_rect(slide, 0, 0, Inches(0.12), H, fill_color=MID_GREEN)
    # タイトル帯
    add_rect(slide, Inches(0.25), Inches(0.2), W - Inches(0.5), Inches(0.9), fill_color=OFF_WHITE)
    # タイトルオレンジ下線
    add_rect(slide, Inches(0.25), Inches(0.95), Inches(6), Inches(0.05), fill_color=ORANGE)
    # タイトルテキスト
    add_textbox(slide, title_text,
                Inches(0.35), Inches(0.2), Inches(10), Inches(0.7),
                font_size=26, bold=True, color=DARK_GREEN)
    # サブタイトル
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.35), Inches(0.95), Inches(10), Inches(0.45),
                    font_size=14, color=MID_GRAY)
    # フッター
    add_textbox(slide, "ステップアップナビ  |  Confidential",
                Inches(8.5), Inches(7.1), Inches(4.5), Inches(0.35),
                font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def kpi_box(slide, label, value, left, top, width=Inches(3.0), height=Inches(1.1)):
    """オレンジ強調ボックス"""
    add_rect(slide, left, top, width, height, fill_color=ORANGE)
    add_textbox(slide, value, left, top + Inches(0.1), width, Inches(0.6),
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + Inches(0.65), width, Inches(0.4),
                font_size=12, color=DARK_GREEN, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows, left, top, width, height,
              header_bg=MID_GREEN, row_alt=MINT, font_size=13):
    col_count = len(headers)
    row_count  = len(rows) + 1
    tbl = slide.shapes.add_table(row_count, col_count, left, top, width, height).table

    col_w = width // col_count
    for i in range(col_count):
        tbl.columns[i].width = col_w

    def style_cell(cell, text, bg, fg=WHITE, bold=False, sz=font_size, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Meiryo"

    for ci, h in enumerate(headers):
        style_cell(tbl.cell(0, ci), h, header_bg, WHITE, bold=True)

    for ri, row in enumerate(rows):
        bg = MINT if ri % 2 == 1 else WHITE
        for ci, val in enumerate(row):
            style_cell(tbl.cell(ri + 1, ci), val, bg, DARK_TEXT, align=PP_ALIGN.CENTER)

    return tbl


# ══════════════════════════════════════════════════════════════
# スライド 1：表紙
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK)
# 全面 DARK_GREEN
add_rect(slide1, 0, 0, W, H, fill_color=DARK_GREEN)
# 下部オレンジ帯
add_rect(slide1, 0, Inches(6.2), W, Inches(1.3), fill_color=ORANGE_DARK)
# 装飾：右上ライトグリーン矩形
add_rect(slide1, Inches(10), 0, Inches(3.33), Inches(3.5), fill_color=MID_GREEN)
# タイトル
add_textbox(slide1, "ステップアップナビ",
            Inches(0.7), Inches(1.5), Inches(9), Inches(1.4),
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
# サブ
add_textbox(slide1, "軽度知的障害者の就労を、テクノロジーで支援する",
            Inches(0.7), Inches(2.9), Inches(10), Inches(0.8),
            font_size=22, color=LIGHT_GREEN, align=PP_ALIGN.LEFT)
# キャッチ（オレンジ帯内）
add_textbox(slide1, "すべての人に、はたらく一歩を。",
            Inches(0.7), Inches(6.25), Inches(9), Inches(0.8),
            font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
# ロゴエリア（右下）
add_textbox(slide1, "PITCH DECK  2026",
            Inches(9.5), Inches(6.3), Inches(3.5), Inches(0.6),
            font_size=14, color=WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# スライド 2：解決する社会課題
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK)
body_slide_base(slide2, "解決する社会課題", "なぜ、いま「ステップアップナビ」が必要なのか")

issues = [
    ("約100万人以上", "日本の軽度知的障害者数（潜在含む）"),
    ("約40万人", "就労系障害福祉サービス利用者（2023年）"),
    ("毎年増加中", "就労継続支援A型・B型の利用者数"),
]
bx = Inches(0.35)
for i, (val, lbl) in enumerate(issues):
    kpi_box(slide2, lbl, val, bx + i * Inches(3.3), Inches(1.35), width=Inches(3.1), height=Inches(1.1))

problems = [
    "「どんな仕事に向いているかわからない」—— 就職先を選ぶ根拠が乏しい",
    "「成長の見通しが立てられない」—— 毎日の積み重ねが見えない",
    "「支援者間で情報が分断される」—— 家族・支援員・企業がバラバラに動く",
    "現状の支援は紙・口頭が中心 —— デジタル化が10年以上遅れている",
]
add_multiline_textbox(slide2, problems,
                      Inches(0.4), Inches(2.65), Inches(12.5), Inches(3.5),
                      font_size=17, color=DARK_TEXT, line_spacing=Pt(30))

add_rect(slide2, Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.6), fill_color=MINT)
add_textbox(slide2, "▶  これらの課題を、ひとつのプラットフォームで解決するのが「ステップアップナビ」です。",
            Inches(0.5), Inches(6.2), Inches(12.5), Inches(0.6),
            font_size=14, bold=True, color=DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# スライド 3：ターゲットユーザー
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK)
body_slide_base(slide3, "ターゲットユーザー", "誰のためのサービスか")

# 主ユーザーボックス
add_rect(slide3, Inches(0.35), Inches(1.3), Inches(5.8), Inches(2.5), fill_color=MINT)
add_textbox(slide3, "主ユーザー（利用者本人）",
            Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.5),
            font_size=16, bold=True, color=DARK_GREEN)
add_multiline_textbox(slide3, [
    "・ 軽度知的障害を持つ 18〜40代の大人",
    "・ 就労継続支援A型・B型の利用者",
    "・ 就労移行支援サービスを探している方",
    "・ ひらがな表示で安心して使えるUI",
], Inches(0.5), Inches(1.9), Inches(5.5), Inches(1.7), font_size=15, color=DARK_TEXT)

# サブユーザーボックス
add_rect(slide3, Inches(6.5), Inches(1.3), Inches(6.4), Inches(2.5), fill_color=RGBColor(0xE8, 0xF4, 0xE8))
add_textbox(slide3, "サブユーザー（支援者・関係者）",
            Inches(6.65), Inches(1.4), Inches(6), Inches(0.5),
            font_size=16, bold=True, color=DARK_GREEN)
add_multiline_textbox(slide3, [
    "・ 家族（保護者）",
    "・ 就労支援員・作業所スタッフ",
    "・ ジョブコーチ",
    "・ 就労移行支援事業所の管理者",
], Inches(6.65), Inches(1.9), Inches(6), Inches(1.7), font_size=15, color=DARK_TEXT)

# 市場規模
add_rect(slide3, Inches(0.35), Inches(4.1), Inches(12.6), Inches(2.7), fill_color=OFF_WHITE)
add_textbox(slide3, "市場規模",
            Inches(0.5), Inches(4.15), Inches(5), Inches(0.45),
            font_size=17, bold=True, color=DARK_GREEN)

kpi_box(slide3, "就労系障害福祉サービス利用者（2023年）", "約 40万人",
        Inches(0.5), Inches(4.6), Inches(3.8), Inches(1.1))
kpi_box(slide3, "就労移行支援市場規模（2023年）", "約 4,000億円",
        Inches(4.7), Inches(4.6), Inches(3.8), Inches(1.1))
kpi_box(slide3, "法定雇用率（2026年目標）", "2.7%",
        Inches(8.9), Inches(4.6), Inches(3.8), Inches(1.1))


# ══════════════════════════════════════════════════════════════
# スライド 4：プロダクト概要
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK)
body_slide_base(slide4, "プロダクト概要", "ステップアップナビとは")

add_rect(slide4, Inches(0.35), Inches(1.3), Inches(12.6), Inches(1.2), fill_color=MINT)
add_textbox(slide4,
            "「特性診断 → 成長ロードマップ → 日々の記録 → 就職マッチング」を一気通貫でサポートするWebアプリ",
            Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.1),
            font_size=17, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

specs = [
    ("アプリ名",   "ステップアップナビ（Step Up Navi）"),
    ("対象",     "軽度知的障害を持つ大人 + 支援者（家族・支援員・事業所スタッフ）"),
    ("技術スタック", "Django（Python）/ PostgreSQL / OpenAI GPT-4o-mini / Railway"),
    ("現在の状況",  "β版 稼働中 — Railway 本番環境にてサービス提供中"),
    ("本番URL",   "https://web-production-c0abe.up.railway.app/"),
    ("ひらがな対応", "全ページでひらがな ↔ 漢字切り替え（セッション保持）"),
]
for i, (k, v) in enumerate(specs):
    y = Inches(2.7) + i * Inches(0.65)
    bg = MINT if i % 2 == 0 else OFF_WHITE
    add_rect(slide4, Inches(0.35), y, Inches(12.6), Inches(0.6), fill_color=bg)
    add_textbox(slide4, k, Inches(0.45), y + Inches(0.08), Inches(2.2), Inches(0.45),
                font_size=13, bold=True, color=MID_GREEN)
    add_textbox(slide4, v, Inches(2.75), y + Inches(0.08), Inches(10.0), Inches(0.45),
                font_size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════
# スライド 5：主要機能（完成形）
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK)
body_slide_base(slide5, "主要機能（完成形）", "6つのコア機能が連動する統合プラットフォーム")

features = [
    ("🔍", "特性診断", "6特性×AIで強み・課題を分析\n13問のやさしい日本語で実施"),
    ("🗺️", "成長ロードマップ", "仕事タイプ別3ステップ\n具体的な行動指針を提示"),
    ("🤖", "AIアバターコーチ", "自分似のアバターとAIチャット\nGPT-4o-miniで個別アドバイス"),
    ("📊", "支援者ダッシュボード", "多角的評価・連携ノート\n支援チーム間で情報共有"),
    ("🏢", "企業マッチング", "地元企業との特性マッチング\n福祉的就労・一般就労に対応"),
    ("👥", "先輩事例DB", "似た特性の先輩の就職先を紹介\n就労の現実的イメージを提供"),
]

cols = 3
for idx, (icon, name, desc) in enumerate(features):
    col = idx % cols
    row = idx // cols
    lft = Inches(0.35) + col * Inches(4.3)
    top = Inches(1.45) + row * Inches(2.5)
    box_w = Inches(4.0)
    box_h = Inches(2.2)

    add_rect(slide5, lft, top, box_w, box_h, fill_color=WHITE,
             line_color=MID_GREEN, line_width=Pt(1.5))
    # アイコン + 名前帯
    add_rect(slide5, lft, top, box_w, Inches(0.65), fill_color=MID_GREEN)
    add_textbox(slide5, f"{icon}  {name}",
                lft + Inches(0.1), top + Inches(0.08), box_w - Inches(0.2), Inches(0.5),
                font_size=16, bold=True, color=WHITE)
    add_multiline_textbox(slide5, desc.split("\n"),
                          lft + Inches(0.15), top + Inches(0.75), box_w - Inches(0.3), Inches(1.3),
                          font_size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════
# スライド 6：技術的差別化（RAG×AI）
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK)
body_slide_base(slide6, "技術的差別化", "RAG × AI で「あなただけのアドバイス」を実現")

# 中央大ボックス
add_rect(slide6, Inches(0.35), Inches(1.3), Inches(12.6), Inches(1.0), fill_color=DARK_GREEN)
add_textbox(slide6,
            "単なる診断ツールではなく、継続的な成長支援プラットフォーム",
            Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.9),
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tech_items = [
    ("個人データ × RAG",
     "診断結果・日々の記録・支援者評価をベクトルDBに蓄積。\n"
     "OpenAI GPT-4o-mini が「あなたの文脈」に合わせた回答を生成。"),
    ("pgvector で低コスト高品質",
     "PostgreSQL 拡張の pgvector を採用。\n"
     "クラウドDBコストを最小化しながら、高品質なセマンティック検索を実現。"),
    ("フォールバック設計",
     "APIキー未設定でもルールベースで動作。\n"
     "どんな事業所でも導入ゼロコストで開始可能。"),
    ("時系列パターン認識",
     "感情ログ・タスク完了率の変化をAIが自動分析。\n"
     "「月曜日に気分が落ちやすい」などの洞察を支援者に提供。"),
]

for i, (title, body) in enumerate(tech_items):
    col = i % 2
    row = i // 2
    lft = Inches(0.35) + col * Inches(6.5)
    top = Inches(2.5) + row * Inches(2.1)
    add_rect(slide6, lft, top, Inches(6.1), Inches(1.9),
             fill_color=MINT, line_color=MID_GREEN, line_width=Pt(1))
    add_textbox(slide6, title,
                lft + Inches(0.15), top + Inches(0.1), Inches(5.8), Inches(0.45),
                font_size=15, bold=True, color=DARK_GREEN)
    add_multiline_textbox(slide6, body.split("\n"),
                          lft + Inches(0.15), top + Inches(0.55), Inches(5.8), Inches(1.2),
                          font_size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════
# スライド 7：ユーザー体験フロー（Before / After）
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK)
body_slide_base(slide7, "ユーザー体験フロー", "Before → After で変わる就労支援体験")

# Before
add_rect(slide7, Inches(0.35), Inches(1.35), Inches(5.5), Inches(5.5), fill_color=RGBColor(0xFF, 0xF0, 0xE0))
add_rect(slide7, Inches(0.35), Inches(1.35), Inches(5.5), Inches(0.55), fill_color=ORANGE)
add_textbox(slide7, "BEFORE（現状の課題）",
            Inches(0.45), Inches(1.4), Inches(5.3), Inches(0.45),
            font_size=16, bold=True, color=WHITE)
befores = [
    "😞  「自分に向いている仕事がわからない」",
    "📞  支援者に毎回口頭で状況を説明する",
    "📝  紙の記録が散在し、引き継ぎが困難",
    "😔  「成長している実感がない」",
    "🔄  家族・支援員・企業が別々に動く",
    "📄  デジタル記録がなく振り返れない",
]
add_multiline_textbox(slide7, befores,
                      Inches(0.5), Inches(2.1), Inches(5.2), Inches(4.5),
                      font_size=15, color=DARK_TEXT, line_spacing=Pt(28))

# 矢印
add_textbox(slide7, "→", Inches(6.0), Inches(3.8), Inches(1.0), Inches(0.8),
            font_size=40, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)

# After
add_rect(slide7, Inches(7.2), Inches(1.35), Inches(5.8), Inches(5.5), fill_color=MINT)
add_rect(slide7, Inches(7.2), Inches(1.35), Inches(5.8), Inches(0.55), fill_color=DARK_GREEN)
add_textbox(slide7, "AFTER（ステップアップナビ）",
            Inches(7.3), Inches(1.4), Inches(5.6), Inches(0.45),
            font_size=16, bold=True, color=WHITE)
afters = [
    "✅  診断で強みが可視化される",
    "📱  支援者と記録をリアルタイム共有",
    "📅  進捗カレンダーで継続が見える化",
    "🌱  毎日のタスクで成長が実感できる",
    "🤝  支援チーム全員が同じ情報を共有",
    "🤖  AIが個別アドバイスを24時間提供",
]
add_multiline_textbox(slide7, afters,
                      Inches(7.35), Inches(2.1), Inches(5.5), Inches(4.5),
                      font_size=15, color=DARK_TEXT, line_spacing=Pt(28))


# ══════════════════════════════════════════════════════════════
# スライド 8：市場機会
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(BLANK)
body_slide_base(slide8, "市場機会", "拡大する障害者就労支援市場")

kpis = [
    ("就労移行支援市場規模\n（2023年）", "約4,000億円"),
    ("法定雇用率の引き上げ\n（2026年）", "2.7%"),
    ("就労継続支援\n利用者数増加率（5年）", "+35%"),
    ("アジア圏就労支援市場\n（推計）", "数兆円規模"),
]
for i, (lbl, val) in enumerate(kpis):
    kpi_box(slide8, lbl, val,
            Inches(0.35) + i * Inches(3.25), Inches(1.35),
            width=Inches(3.0), height=Inches(1.3))

opportunities = [
    "📈  障害者雇用促進法改正 — 企業の法定雇用率が段階的に引き上げ（2.5%→2.7%）。\n    企業側の障害者採用ニーズが急増。",
    "🏛️  行政のデジタル化推進 — 障害福祉サービスのDX支援が国・自治体レベルで加速。\n    補助金・実証事業の機会が拡大中。",
    "🌏  海外展開可能性 — アジア圏（韓国・台湾・タイ等）でも知的障害者就労支援の\n    デジタル化ニーズは高く、日本モデルの横展開が可能。",
    "🤝  企業の ESG・D&I 意識の高まり — 障害者雇用をブランド価値につなげたい\n    企業とのマッチング需要が増加。",
]
add_multiline_textbox(slide8, opportunities,
                      Inches(0.4), Inches(2.9), Inches(12.5), Inches(3.8),
                      font_size=15, color=DARK_TEXT, line_spacing=Pt(26))


# ══════════════════════════════════════════════════════════════
# スライド 9：ビジネスモデル
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(BLANK)
body_slide_base(slide9, "ビジネスモデル", "3つの収益源で安定・拡張性のある事業設計")

plans = [
    ("💰 SaaS月額課金（事業所向け）", MID_GREEN, [
        "基本プラン（20名まで）：¥9,800 / 月",
        "スタンダード（100名まで）：¥29,800 / 月",
        "エンタープライズ：要相談",
        "▶ 安定した月次収益（MRR）の柱",
    ]),
    ("🤝 企業マッチング手数料", ORANGE_DARK, [
        "採用成立時に成功報酬",
        "¥50,000 〜 ¥100,000 / 件",
        "福祉的就労・一般就労の両対応",
        "▶ 企業・利用者・事業所の三方よし",
    ]),
    ("📊 データ分析レポート（行政・研究機関向け）", DARK_GREEN, [
        "匿名化した就労成功データの提供",
        "厚生労働省・自治体・大学との連携",
        "政策立案・福祉研究への貢献",
        "▶ 社会的信頼性の向上にも寄与",
    ]),
]

for i, (title, bg, items) in enumerate(plans):
    lft = Inches(0.35) + i * Inches(4.3)
    top = Inches(1.35)
    w   = Inches(4.0)
    add_rect(slide9, lft, top, w, Inches(0.6), fill_color=bg)
    add_textbox(slide9, title, lft + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.5),
                font_size=13, bold=True, color=WHITE)
    body_top = top + Inches(0.65)
    add_rect(slide9, lft, body_top, w, Inches(5.05), fill_color=MINT)
    add_multiline_textbox(slide9, items,
                          lft + Inches(0.15), body_top + Inches(0.15), w - Inches(0.3), Inches(4.7),
                          font_size=14, color=DARK_TEXT, line_spacing=Pt(26))


# ══════════════════════════════════════════════════════════════
# スライド 10：収益シミュレーション（3年計画）
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(BLANK)
body_slide_base(slide10, "収益シミュレーション（3年計画）", "段階的な成長戦略と売上見通し")

headers = ["", "1年目", "2年目", "3年目"]
rows = [
    ["導入事業所数",      "10所",      "50所",      "200所"],
    ["SaaS月額収入",     "98万円/月",  "490万円/月", "1,960万円/月"],
    ["マッチング収入",    "50万円/年",  "500万円/年", "3,000万円/年"],
    ["年間売上（合計）", "約1,200万円", "約6,400万円", "約2.65億円"],
]
add_table(slide10, headers, rows,
          Inches(0.35), Inches(1.4), Inches(12.6), Inches(2.5),
          font_size=15)

# KPI 強調
kpis10 = [
    ("3年目 年間売上目標", "約2.65億円"),
    ("3年目 導入事業所", "200所"),
    ("事業所あたり\n平均月額収入", "98万円〜"),
]
for i, (lbl, val) in enumerate(kpis10):
    kpi_box(slide10, lbl, val,
            Inches(0.35) + i * Inches(4.3), Inches(4.25),
            width=Inches(4.0), height=Inches(1.2))

add_rect(slide10, Inches(0.35), Inches(5.65), Inches(12.6), Inches(0.7), fill_color=MINT)
add_textbox(slide10,
            "※ 収益は SaaS 月額 × 導入事業所数 + マッチング成功報酬の積み上げ。エンタープライズ契約・行政案件は含まず。",
            Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════
# スライド 11：競合比較
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(BLANK)
body_slide_base(slide11, "競合比較", "ステップアップナビが持つ唯一無二の強み")

headers11 = ["比較項目", "ステップアップナビ", "既存支援ツール\n（紙・Excel）", "一般就活アプリ"]
rows11 = [
    ["知的障害者特化",    "✅",         "△",    "❌"],
    ["AI個別アドバイス",  "✅",         "❌",    "△"],
    ["支援者連携",       "✅",         "△",    "❌"],
    ["企業マッチング",    "✅",         "❌",    "✅"],
    ["ひらがな対応",      "✅",         "❌",    "❌"],
    ["成長ロードマップ",   "✅",         "❌",    "❌"],
    ["月額コスト（事業所）", "低（SaaS）", "0（非効率）", "中〜高"],
]
add_table(slide11, headers11, rows11,
          Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.5),
          font_size=14)


# ══════════════════════════════════════════════════════════════
# スライド 12：開発ロードマップ（資金使途と紐付け）
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(BLANK)
body_slide_base(slide12, "開発ロードマップ", "フェーズ別の資金使途と機能展開計画")

phases = [
    ("現在\n稼働中", DARK_GREEN, [
        "✅ 特性診断（6特性・13問）",
        "✅ 成長ロードマップ（3ステップ）",
        "✅ AIアバターチャット",
        "✅ 支援者ダッシュボード",
        "✅ Railway本番稼働",
    ], "β版稼働中"),
    ("フェーズ2\n〜6ヶ月", MID_GREEN, [
        "◎ 支援者評価レーダーチャート",
        "◎ 進捗カレンダー・ストリーク",
        "◎ 支援者間連携ノート",
        "◎ 先輩事例データベース",
        "◎ ルーティンチェックリスト",
    ], "必要資金：500万円"),
    ("フェーズ3\n〜12ヶ月", ORANGE_DARK, [
        "△ RAG実装（pgvector）",
        "△ AI時系列パターン分析",
        "△ 職場ロールプレイ練習",
        "◇ 企業マッチング機能",
        "◇ 多職種連携プラットフォーム",
    ], "必要資金：1,500万円"),
    ("フェーズ4\n〜24ヶ月", RGBColor(0x8B, 0x45, 0x13), [
        "◇ 地元企業タイアップ展開",
        "◇ スマートフォンアプリ化",
        "◇ アジア圏への海外展開",
        "◇ 行政・研究機関データ提供",
        "◇ 多言語対応",
    ], "必要資金：5,000万円"),
]

for i, (phase, bg, items, budget) in enumerate(phases):
    lft = Inches(0.35) + i * Inches(3.25)
    top = Inches(1.35)
    w   = Inches(3.0)

    add_rect(slide12, lft, top, w, Inches(0.75), fill_color=bg)
    add_textbox(slide12, phase, lft + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.65),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    body_top = top + Inches(0.8)
    add_rect(slide12, lft, body_top, w, Inches(4.3), fill_color=OFF_WHITE,
             line_color=bg, line_width=Pt(1.5))
    add_multiline_textbox(slide12, items,
                          lft + Inches(0.1), body_top + Inches(0.1), w - Inches(0.2), Inches(3.5),
                          font_size=12, color=DARK_TEXT, line_spacing=Pt(22))

    add_rect(slide12, lft, body_top + Inches(3.7), w, Inches(0.5), fill_color=bg)
    add_textbox(slide12, budget,
                lft + Inches(0.05), body_top + Inches(3.73), w - Inches(0.1), Inches(0.45),
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 13：チーム・開発体制
# ══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(BLANK)
body_slide_base(slide13, "チーム・開発体制", "現場の声を反映し続ける開発チーム")

members = [
    ("👤", "プロダクトオーナー", [
        "就労支援現場 10年以上の経験",
        "利用者・支援者の声を直接収集",
        "プロダクト戦略・事業開発担当",
    ]),
    ("💻", "テックリード", [
        "Python / Django / AI実装",
        "OpenAI API × pgvector 設計",
        "Railway インフラ・CI/CD 管理",
    ]),
    ("🏥", "福祉専門家アドバイザー", [
        "社会福祉士・精神保健福祉士",
        "就労支援プログラムの監修",
        "ユーザーテスト・フィードバック収集",
    ]),
    ("📊", "ビジネス開発", [
        "事業所・企業との連携交渉",
        "行政・補助金申請サポート",
        "マーケティング・PR戦略",
    ]),
]

for i, (icon, role, items) in enumerate(members):
    col = i % 2
    row = i // 2
    lft = Inches(0.35) + col * Inches(6.5)
    top = Inches(1.4) + row * Inches(2.6)
    w = Inches(6.1)
    h = Inches(2.4)
    add_rect(slide13, lft, top, w, h, fill_color=MINT,
             line_color=MID_GREEN, line_width=Pt(1))
    add_textbox(slide13, f"{icon}  {role}",
                lft + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.5),
                font_size=17, bold=True, color=DARK_GREEN)
    add_multiline_textbox(slide13, [f"・ {it}" for it in items],
                          lft + Inches(0.15), top + Inches(0.65), w - Inches(0.3), Inches(1.6),
                          font_size=14, color=DARK_TEXT, line_spacing=Pt(22))

add_rect(slide13, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.5), fill_color=DARK_GREEN)
add_textbox(slide13, "「障害のある方の現場経験」と「最先端AI技術」を掛け合わせた開発体制",
            Inches(0.5), Inches(6.73), Inches(12.3), Inches(0.45),
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 14：求める支援
# ══════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(BLANK)
body_slide_base(slide14, "求める支援", "シードラウンドで実現すること")

# 資金調達ボックス
add_rect(slide14, Inches(0.35), Inches(1.35), Inches(12.6), Inches(1.3), fill_color=DARK_GREEN)
add_textbox(slide14, "資金調達目標：シードラウンド  3,000万円",
            Inches(0.5), Inches(1.45), Inches(7), Inches(0.7),
            font_size=24, bold=True, color=WHITE)
add_textbox(slide14, "（エクイティ / 補助金 / 融資の組み合わせ対応可）",
            Inches(0.5), Inches(2.0), Inches(7), Inches(0.55),
            font_size=13, color=LIGHT_GREEN)

kpi_box(slide14, "シードラウンド目標", "3,000万円",
        Inches(9.0), Inches(1.35), Inches(3.9), Inches(1.3))

uses = [
    ("開発人員強化（50%）", "エンジニア2名採用・フリーランス連携強化"),
    ("事業所トライアル提供（20%）", "初期50事業所に3ヶ月無料導入・改善フィードバック収集"),
    ("マーケティング（15%）", "福祉業界メディア・展示会への出展・PR活動"),
    ("インフラ・AIコスト（10%）", "Railway本番スケール・OpenAI API費用"),
    ("法務・コンプライアンス（5%）", "個人情報保護・障害福祉関連法規対応"),
]
add_textbox(slide14, "【資金使途】",
            Inches(0.4), Inches(2.85), Inches(5), Inches(0.4),
            font_size=14, bold=True, color=DARK_GREEN)
for i, (k, v) in enumerate(uses):
    y = Inches(3.25) + i * Inches(0.55)
    bg = MINT if i % 2 == 0 else OFF_WHITE
    add_rect(slide14, Inches(0.35), y, Inches(12.6), Inches(0.5), fill_color=bg)
    add_textbox(slide14, k, Inches(0.45), y + Inches(0.07), Inches(3.2), Inches(0.4),
                font_size=12, bold=True, color=MID_GREEN)
    add_textbox(slide14, v, Inches(3.8), y + Inches(0.07), Inches(9), Inches(0.4),
                font_size=12, color=DARK_TEXT)

add_rect(slide14, Inches(0.35), Inches(6.1), Inches(12.6), Inches(0.8), fill_color=MINT)
add_textbox(slide14,
            "資金以外に求めるもの：  就労支援事業所とのネットワーク  ／  地元企業との連携  ／  福祉・医療の専門知識・人脈",
            Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.65),
            font_size=14, bold=True, color=DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# スライド 15：まとめ・ビジョン
# ══════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(BLANK)
add_rect(slide15, 0, 0, W, H, fill_color=DARK_GREEN)
# 右上装飾
add_rect(slide15, Inches(10), 0, Inches(3.33), Inches(4), fill_color=MID_GREEN)
# 下オレンジ帯
add_rect(slide15, 0, Inches(5.8), W, Inches(1.7), fill_color=ORANGE_DARK)

add_textbox(slide15, "ビジョン",
            Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
            font_size=18, color=LIGHT_GREEN)
add_textbox(slide15, "障害のある方が、\n自分らしくはたらける社会へ",
            Inches(0.7), Inches(1.0), Inches(9.2), Inches(1.8),
            font_size=38, bold=True, color=WHITE)

milestones = [
    ("3年後", "全国の就労支援事業所の 10% に導入\n年間売上 約2.65億円"),
    ("5年後", "アジア圏（韓国・台湾・タイ等）への海外展開\n利用者 10万人突破"),
]
for i, (year, desc) in enumerate(milestones):
    lft = Inches(0.7) + i * Inches(6.0)
    top = Inches(3.2)
    add_rect(slide15, lft, top, Inches(5.5), Inches(2.0), fill_color=MID_GREEN)
    add_textbox(slide15, year,
                lft + Inches(0.2), top + Inches(0.1), Inches(5.1), Inches(0.45),
                font_size=18, bold=True, color=ORANGE)
    add_multiline_textbox(slide15, desc.split("\n"),
                          lft + Inches(0.2), top + Inches(0.6), Inches(5.1), Inches(1.3),
                          font_size=16, color=WHITE)

add_textbox(slide15,
            "一人ひとりの『はたらく一歩』を、テクノロジーで支える",
            Inches(0.7), Inches(5.9), Inches(11.5), Inches(0.7),
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide15,
            "ステップアップナビ  |  https://web-production-c0abe.up.railway.app/",
            Inches(0.7), Inches(6.65), Inches(11.5), Inches(0.5),
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ─── 保存 ────────────────────────────────────────────────────
OUTPUT = "/Users/matsunaganaoto/Desktop/projects/App/ステップアップナビ_ピッチデック.pptx"
prs.save(OUTPUT)

slide_count = len(prs.slides)
print(f"✅ ピッチデック作成完了")
print(f"   ファイル：{OUTPUT}")
print(f"   スライド数：{slide_count}枚")
