from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ─────────────────────────────────────────────
# カラーパレット（新デザインシステム）
# ─────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x1B, 0x45, 0x32)
MID_GREEN   = RGBColor(0x2D, 0x6A, 0x4F)
LIGHT_GREEN = RGBColor(0x95, 0xD5, 0xB2)
MINT        = RGBColor(0xD8, 0xF3, 0xDC)

ORANGE      = RGBColor(0xF4, 0xA2, 0x61)
ORANGE_DARK = RGBColor(0xE7, 0x6F, 0x51)

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xF8)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY    = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF0)
SHADOW_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

PHASE2_BLUE = RGBColor(0x23, 0x7A, 0xB0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_vertical_accent(slide):
    """左端に細い縦ライン（MID_GREEN アクセント）"""
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(0.08), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_GREEN
    bar.line.fill.background()


def add_footer(slide):
    """右下にフッターテキスト"""
    add_text_box(slide, "🌱 ステップアップナビ",
                 Inches(10.8), Inches(7.1),
                 Inches(2.4), Inches(0.35),
                 font_size=8, color=MID_GRAY,
                 align=PP_ALIGN.RIGHT)


def add_slide_decorations(slide):
    """左縦ライン＋フッターを一括追加"""
    add_vertical_accent(slide)
    add_footer(slide)


def add_slide_title(slide, title, subtitle=None):
    """新デザインのスライドタイトル"""
    add_text_box(slide, title,
                 Inches(0.3), Inches(0.3),
                 Inches(12.5), Inches(0.7),
                 font_size=28, bold=True, color=DARK_TEXT)

    # オレンジの下線
    underline = slide.shapes.add_shape(
        1,
        Inches(0.3), Inches(1.02),
        Inches(3.0), Inches(0.04)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(1.1),
                     Inches(12.5), Inches(0.4),
                     font_size=13, color=MID_GRAY)


def add_card_shadow(slide, x, y, w, h):
    """影風の薄グレー矩形"""
    shadow = slide.shapes.add_shape(
        1,
        x + Inches(0.06), y + Inches(0.06),
        w, h
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SHADOW_GRAY
    shadow.line.fill.background()


def add_card(slide, x, y, w, h, border_color=None):
    """白背景カード＋左ボーダー"""
    add_card_shadow(slide, x, y, w, h)
    card = slide.shapes.add_shape(1, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    bc = border_color or MID_GREEN
    left_bar = slide.shapes.add_shape(
        1, x, y, Inches(0.1), h
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = bc
    left_bar.line.fill.background()
    return card


# ─────────────────────────────────────────────
# スライド1：表紙
# ─────────────────────────────────────────────
def add_title_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, DARK_GREEN)

    # 下部オレンジ帯
    orange_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.35),
        SLIDE_W, Inches(0.15)
    )
    orange_bar.fill.solid()
    orange_bar.fill.fore_color.rgb = ORANGE
    orange_bar.line.fill.background()

    # 右下装飾円
    circle = slide.shapes.add_shape(9, Inches(10.0), Inches(4.5), Inches(4.5), Inches(4.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x15, 0x3A, 0x28)
    circle.line.fill.background()

    # 左上装飾円
    circle2 = slide.shapes.add_shape(9, Inches(-1.2), Inches(-1.2), Inches(3.5), Inches(3.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x15, 0x3A, 0x28)
    circle2.line.fill.background()

    # 絵文字
    add_text_box(slide, "🌱",
                 Inches(5.5), Inches(1.3),
                 Inches(2.3), Inches(1.1),
                 font_size=48, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # メインタイトル
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ステップアップナビ"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # サブタイトル
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.85), Inches(11.3), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "就労支援Webアプリ　機能紹介・開発ロードマップ"
    run2.font.size = Pt(20)
    run2.font.color.rgb = LIGHT_GREEN

    # URL・バージョン日付
    add_text_box(slide, "https://web-production-c0abe.up.railway.app/　｜　2026年6月27日",
                 Inches(7.5), Inches(6.9),
                 Inches(5.5), Inches(0.4),
                 font_size=10, color=WHITE,
                 align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# スライド2：アプリ概要
# ─────────────────────────────────────────────
def add_overview_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "アプリの目的")

    items = [
        ("🎯  対象", "軽度知的障害を持つ大人"),
        ("🌱  目的", "就労（福祉的就労・一般就労）に向けた成長を支援するWebアプリ"),
        ("👥  利用者", "本人（障害のある大人）＋ 支援者（家族・支援員・作業所スタッフ）"),
        ("💻  技術スタック", "Django（Python）、PostgreSQL、Railway（本番環境）"),
        ("🔗  本番URL", "https://web-production-c0abe.up.railway.app/"),
    ]

    y = Inches(1.6)
    for label, content in items:
        card_h = Inches(0.78)
        add_card_shadow(slide, Inches(0.25), y, Inches(12.8), card_h)

        card = slide.shapes.add_shape(1, Inches(0.25), y, Inches(12.8), card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        # 左ボーダー
        lb = slide.shapes.add_shape(1, Inches(0.25), y, Inches(0.1), card_h)
        lb.fill.solid()
        lb.fill.fore_color.rgb = MID_GREEN
        lb.line.fill.background()

        # ラベル背景
        lbl_bg = slide.shapes.add_shape(1, Inches(0.35), y, Inches(2.6), card_h)
        lbl_bg.fill.solid()
        lbl_bg.fill.fore_color.rgb = MID_GREEN
        lbl_bg.line.fill.background()

        add_text_box(slide, label,
                     Inches(0.42), y + Inches(0.1),
                     Inches(2.45), Inches(0.58),
                     font_size=15, bold=True, color=WHITE)

        add_text_box(slide, content,
                     Inches(3.15), y + Inches(0.12),
                     Inches(9.7), Inches(0.55),
                     font_size=16, color=DARK_TEXT)

        y += Inches(0.95)


# ─────────────────────────────────────────────
# スライド3：現在の実装機能（概要）
# ─────────────────────────────────────────────
def add_current_features_overview(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "★ 現在すでに動いている機能（9機能）",
                    "すべてRailway本番環境で稼働中")

    features = [
        ("🔐 利用者登録・ログイン", "アカウント登録・支援者紐付け・ログイン/アウト"),
        ("🧠 性格・特性診断", "13問・6特性・5仕事タイプ判定"),
        ("🗺️ 成長ロードマップ", "仕事タイプ別3ステップで構成"),
        ("📓 日々の記録・ダッシュボード", "感情スタンプ・体調・週間グラフ"),
        ("🧑 自分似のアバター機能", "SVGカスタマイズ・診断連動"),
        ("💬 アバターとのAIチャット", "GPT-4o-mini / フォールバック対応"),
        ("👩‍💼 支援者ダッシュボード", "進捗確認・メモ機能"),
        ("🈶 漢字/ひらがな切り替え", "全ページ対応・セッション保持"),
        ("🚀 本番デプロイ環境", "Railway + PostgreSQL で稼働中"),
    ]

    col_w = Inches(4.15)
    cols = [Inches(0.25), Inches(4.57), Inches(8.9)]
    rows_per_col = 3
    y_start = Inches(1.6)
    row_h = Inches(1.6)

    for i, (name, desc) in enumerate(features):
        col_i = i // rows_per_col
        row_i = i % rows_per_col
        x = cols[col_i]
        y = y_start + row_i * row_h
        h = Inches(1.45)

        add_card_shadow(slide, x, y, col_w, h)

        card = slide.shapes.add_shape(1, x, y, col_w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        lb = slide.shapes.add_shape(1, x, y, Inches(0.1), h)
        lb.fill.solid()
        lb.fill.fore_color.rgb = MID_GREEN
        lb.line.fill.background()

        add_text_box(slide, name,
                     x + Inches(0.18), y + Inches(0.1),
                     col_w - Inches(0.28), Inches(0.5),
                     font_size=14, bold=True, color=MID_GREEN)
        add_text_box(slide, desc,
                     x + Inches(0.18), y + Inches(0.62),
                     col_w - Inches(0.28), Inches(0.7),
                     font_size=12, color=DARK_TEXT)


# ─────────────────────────────────────────────
# スライド4：性格・特性診断
# ─────────────────────────────────────────────
def add_diagnosis_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "🧠 性格・特性診断",
                    "実装済み ★ ｜ 13問のやさしい日本語による質問")

    traits = [
        ("集中力", "一つの作業に取り組み続ける力"),
        ("コミュニケーション力", "人と話したり伝えたりする力"),
        ("体力", "身体的な活動への耐久力"),
        ("几帳面さ", "丁寧・正確に作業する力"),
        ("感情コントロール", "気持ちを安定させる力"),
        ("学習意欲", "新しいことを学ぼうとする力"),
    ]

    add_text_box(slide, "📊 測定する6つの特性",
                 Inches(0.3), Inches(1.6),
                 Inches(6.5), Inches(0.5),
                 font_size=17, bold=True, color=DARK_TEXT)

    for i, (trait, desc) in enumerate(traits):
        row = i % 3
        col = i // 3
        x = Inches(0.3) + col * Inches(3.3)
        y = Inches(2.1) + row * Inches(1.3)
        h = Inches(1.15)

        add_card_shadow(slide, x, y, Inches(3.1), h)
        card = slide.shapes.add_shape(1, x, y, Inches(3.1), h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        bar = slide.shapes.add_shape(1, x, y, Inches(0.1), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.fill.background()

        add_text_box(slide, trait,
                     x + Inches(0.18), y + Inches(0.08),
                     Inches(2.8), Inches(0.45),
                     font_size=15, bold=True, color=DARK_TEXT)
        add_text_box(slide, desc,
                     x + Inches(0.18), y + Inches(0.52),
                     Inches(2.8), Inches(0.5),
                     font_size=11, color=MID_GRAY)

    # 5仕事タイプ
    add_text_box(slide, "💼 判定される5つの仕事タイプ",
                 Inches(7.0), Inches(1.6),
                 Inches(6.0), Inches(0.5),
                 font_size=17, bold=True, color=DARK_TEXT)

    job_types = ["🌾 農業", "🏭 製造", "🧹 清掃", "🍱 食品加工", "🤝 接客"]
    for i, jt in enumerate(job_types):
        x = Inches(7.1) + (i % 3) * Inches(2.05)
        y = Inches(2.2) + (i // 3) * Inches(1.0)
        chip = slide.shapes.add_shape(1, x, y, Inches(1.85), Inches(0.75))
        chip.fill.solid()
        chip.fill.fore_color.rgb = MID_GREEN
        chip.line.fill.background()
        add_text_box(slide, jt,
                     x + Inches(0.1), y + Inches(0.1),
                     Inches(1.65), Inches(0.5),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

    add_text_box(slide, "▶ 結果として「強み・課題・仕事タイプ」をわかりやすく表示",
                 Inches(0.3), Inches(6.4),
                 Inches(12.7), Inches(0.55),
                 font_size=15, color=MID_GREEN, bold=True)


# ─────────────────────────────────────────────
# スライド5：成長ロードマップ
# ─────────────────────────────────────────────
def add_roadmap_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "🗺️ 成長ロードマップ",
                    "実装済み ★ ｜ 仕事タイプ別 × 3ステップ構成")

    steps = [
        ("STEP 1", "生活習慣の安定",
         "▶ 規則正しい生活リズムの確立\n▶ 体調管理の基本スキル\n▶ 日々の記録をつける習慣"),
        ("STEP 2", "作業スキルの習得",
         "▶ 仕事タイプ別の基礎作業練習\n▶ 集中力・体力の強化\n▶ 職場でのマナー学習"),
        ("STEP 3", "就労準備",
         "▶ 実際の就労先の見学・体験\n▶ 面接・職場コミュニケーション練習\n▶ 就労後のフォローアップ"),
    ]

    step_colors = [MID_GREEN, PHASE2_BLUE, ORANGE_DARK]

    for i, (step, title, content) in enumerate(steps):
        x = Inches(0.3) + i * Inches(4.32)
        y = Inches(1.55)
        w = Inches(4.1)
        h = Inches(5.5)

        add_card_shadow(slide, x, y, w, h)

        # カード背景
        card = slide.shapes.add_shape(1, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        # ヘッダー帯
        hdr = slide.shapes.add_shape(1, x, y, w, Inches(0.4))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = step_colors[i]
        hdr.line.fill.background()

        add_text_box(slide, step,
                     x + Inches(0.1), y + Inches(0.02),
                     w - Inches(0.2), Inches(0.38),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        # タイトル帯
        add_text_box(slide, title,
                     x + Inches(0.15), y + Inches(0.5),
                     w - Inches(0.3), Inches(0.7),
                     font_size=20, bold=True, color=step_colors[i],
                     align=PP_ALIGN.CENTER)

        # 区切り線
        sep = slide.shapes.add_shape(1, x + Inches(0.3), y + Inches(1.28),
                                     w - Inches(0.6), Inches(0.04))
        sep.fill.solid()
        sep.fill.fore_color.rgb = LIGHT_GREEN
        sep.line.fill.background()

        add_text_box(slide, content,
                     x + Inches(0.2), y + Inches(1.42),
                     w - Inches(0.4), Inches(3.8),
                     font_size=14, color=DARK_TEXT)

        if i < 2:
            arr = slide.shapes.add_shape(
                13,
                x + w + Inches(0.04), Inches(4.0),
                Inches(0.2), Inches(0.5)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    add_text_box(slide, "💡 タブ切り替えUIで仕事タイプ別に見やすく表示",
                 Inches(0.3), Inches(7.1),
                 Inches(12.7), Inches(0.35),
                 font_size=13, color=MID_GRAY)


# ─────────────────────────────────────────────
# スライド6：アバター機能
# ─────────────────────────────────────────────
def add_avatar_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "🧑 自分似のアバター機能",
                    "実装済み ★ ｜ SVGカスタマイズ・診断連動")

    features = [
        ("🎨", "外見カスタマイズ", "肌の色・髪型・目・表情などをSVGで\n自由にカスタマイズ", MID_GREEN),
        ("👕", "仕事タイプ連動制服", "診断結果の仕事タイプに応じた\n制服が自動的に設定される", PHASE2_BLUE),
        ("😊", "特性スコア連動", "特性スコアに応じてアクセサリー・\n表情が自動的に変化", ORANGE_DARK),
        ("📊", "常時表示", "ダッシュボードでアバターを\n常に表示し、成長を実感", MID_GREEN),
    ]

    for i, (emoji, title, desc, accent) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.25) + col * Inches(6.52)
        y = Inches(1.6) + row * Inches(2.55)
        w = Inches(6.3)
        h = Inches(2.35)

        add_card_shadow(slide, x, y, w, h)

        card = slide.shapes.add_shape(1, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        lb = slide.shapes.add_shape(1, x, y, Inches(0.1), h)
        lb.fill.solid()
        lb.fill.fore_color.rgb = accent
        lb.line.fill.background()

        hdr_band = slide.shapes.add_shape(1, x + Inches(0.1), y, w - Inches(0.1), Inches(0.4))
        hdr_band.fill.solid()
        hdr_band.fill.fore_color.rgb = accent
        hdr_band.line.fill.background()

        add_text_box(slide, f"{emoji}  {title}",
                     x + Inches(0.25), y + Inches(0.04),
                     w - Inches(0.4), Inches(0.34),
                     font_size=15, bold=True, color=WHITE)

        add_text_box(slide, desc,
                     x + Inches(0.25), y + Inches(0.5),
                     w - Inches(0.4), Inches(1.7),
                     font_size=13, color=DARK_TEXT)


# ─────────────────────────────────────────────
# スライド7：AIチャット
# ─────────────────────────────────────────────
def add_ai_chat_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "💬 アバターとのAIチャット",
                    "実装済み ★ ｜ GPT-4o-mini搭載・フォールバック対応")

    # 左：API有り
    lx, ly, lw, lh = Inches(0.25), Inches(1.6), Inches(6.3), Inches(5.5)
    add_card_shadow(slide, lx, ly, lw, lh)
    box1 = slide.shapes.add_shape(1, lx, ly, lw, lh)
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.fill.background()

    hdr1 = slide.shapes.add_shape(1, lx, ly, lw, Inches(0.45))
    hdr1.fill.solid()
    hdr1.fill.fore_color.rgb = MID_GREEN
    hdr1.line.fill.background()

    add_text_box(slide, "🤖  OpenAI APIキーあり",
                 lx + Inches(0.15), ly + Inches(0.06),
                 lw - Inches(0.3), Inches(0.36),
                 font_size=16, bold=True, color=WHITE)

    content1 = "▶ GPT-4o-mini による個別対応\n▶ 本人の診断結果・特性を踏まえた会話\n▶ 自然で共感的な対話が可能\n▶ 会話履歴の保持・リセット機能"
    add_text_box(slide, content1,
                 lx + Inches(0.2), ly + Inches(0.6),
                 lw - Inches(0.35), Inches(4.7),
                 font_size=15, color=DARK_TEXT)

    # 右：API無し
    rx, ry, rw, rh = Inches(6.8), Inches(1.6), Inches(6.3), Inches(5.5)
    add_card_shadow(slide, rx, ry, rw, rh)
    box2 = slide.shapes.add_shape(1, rx, ry, rw, rh)
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.fill.background()

    hdr2 = slide.shapes.add_shape(1, rx, ry, rw, Inches(0.45))
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = MID_GRAY
    hdr2.line.fill.background()

    add_text_box(slide, "💡  APIキーなし（フォールバック）",
                 rx + Inches(0.15), ry + Inches(0.06),
                 rw - Inches(0.3), Inches(0.36),
                 font_size=16, bold=True, color=WHITE)

    content2 = "▶ ルールベースの会話で動作\n▶ 誰でも追加費用なしで利用可能\n▶ 基本的な励まし・アドバイスに対応\n▶ APIキー設定でいつでも高機能化可能"
    add_text_box(slide, content2,
                 rx + Inches(0.2), ry + Inches(0.6),
                 rw - Inches(0.35), Inches(4.7),
                 font_size=15, color=DARK_TEXT)

    add_text_box(slide, "⇄",
                 Inches(6.35), Inches(4.1),
                 Inches(0.6), Inches(0.6),
                 font_size=26, color=ORANGE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# スライド8：漢字切り替え・支援者ダッシュボード
# ─────────────────────────────────────────────
def add_accessibility_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "🈶 漢字/ひらがな切り替え ＋ 支援者ダッシュボード",
                    "実装済み ★ ｜ ユニバーサルデザインへの配慮")

    for col_i, (box_title, features, accent) in enumerate([
        ("🈶 漢字 / ひらがな 切り替え機能",
         [
             "✅ 全ページで漢字表記 ↔ ひらがな表記を切り替え",
             "✅ セッションで設定を保持（ページ移動しても維持）",
             "✅ 知的障害のある利用者への読みやすさへの配慮",
             "",
             "📌 ユニバーサルデザインの重要な実装の一つ",
             "　　本人が自分に合った表示で操作できる",
         ], MID_GREEN),
        ("👩‍💼 支援者ダッシュボード",
         [
             "✅ 担当利用者の進捗・記録を一覧確認",
             "✅ 利用者へのメモ機能（支援コメント）",
             "",
             "👨‍👩‍👧  支援者（家族・支援員・作業所スタッフ）が",
             "　　利用者の状況をリアルタイムで把握",
             "",
             "📌 支援者も安心して使えるシンプルなUI",
         ], PHASE2_BLUE),
    ]):
        bx = Inches(0.25) + col_i * Inches(6.6)
        by, bw, bh = Inches(1.6), Inches(6.3), Inches(5.5)

        add_card_shadow(slide, bx, by, bw, bh)
        box = slide.shapes.add_shape(1, bx, by, bw, bh)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.fill.background()

        hdr = slide.shapes.add_shape(1, bx, by, bw, Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = accent
        hdr.line.fill.background()

        add_text_box(slide, box_title,
                     bx + Inches(0.15), by + Inches(0.06),
                     bw - Inches(0.3), Inches(0.36),
                     font_size=15, bold=True, color=WHITE)

        y = by + Inches(0.6)
        for feat in features:
            add_text_box(slide, feat,
                         bx + Inches(0.2), y,
                         bw - Inches(0.35), Inches(0.5),
                         font_size=13, color=DARK_TEXT)
            y += Inches(0.57)


# ─────────────────────────────────────────────
# スライド9：AIなしで追加できる機能
# ─────────────────────────────────────────────
def add_no_ai_features_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "◎ 次のステップで追加できる機能（AIなし）",
                    "既存のDjangoアプリに追加実装するだけで対応可能")

    rows = [
        ("1", "支援者評価 × レーダーチャート", "支援者が6特性を評点入力・自己評価と重ねて表示", "認識ギャップの可視化", "低〜中"),
        ("2", "希望と適性のギャップ表示", "希望タイプと診断結果のズレを表示・次のステップを提示", "目標設定のサポート", "低"),
        ("3", "進捗カレンダー + ストリーク", "月間カレンダーで記録日を色付け・連続記録日数を表示", "習慣化の動機付け", "低"),
        ("4", "朝・夕のルーティンチェック", "仕事前後のやること一覧をチェック形式で管理", "生活習慣の定着", "低"),
        ("5", "先輩事例データベース", "同じ仕事タイプを目指した先輩の体験談を表示", "就労イメージの形成", "中"),
        ("6", "支援者間の連携ノート", "家族・支援員・スタッフが情報共有", "支援の一貫性向上", "中"),
        ("7", "過去診断との比較グラフ", "複数回の診断結果を時系列で比較・成長を可視化", "成長の見える化", "低〜中"),
    ]

    headers = ["#", "機能名", "概要", "期待効果", "難易度"]
    col_xs = [Inches(0.25), Inches(0.75), Inches(3.5), Inches(8.3), Inches(11.5)]
    col_ws = [Inches(0.45), Inches(2.65), Inches(4.65), Inches(3.1), Inches(1.6)]

    hdr_y = Inches(1.6)
    hdr_bg = slide.shapes.add_shape(1, Inches(0.25), hdr_y, Inches(12.85), Inches(0.48))
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = MID_GREEN
    hdr_bg.line.fill.background()

    for j, (h, x, w) in enumerate(zip(headers, col_xs, col_ws)):
        add_text_box(slide, h, x, hdr_y + Inches(0.06), w, Inches(0.38),
                     font_size=13, bold=True, color=WHITE)

    for i, row in enumerate(rows):
        y = Inches(2.13) + i * Inches(0.68)
        bg_color = WHITE if i % 2 == 0 else MINT
        bg = slide.shapes.add_shape(1, Inches(0.25), y - Inches(0.04),
                                    Inches(12.85), Inches(0.65))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.color.rgb = LIGHT_GREEN
        bg.line.width = Pt(0.5)

        for j, (val, x, w) in enumerate(zip(row, col_xs, col_ws)):
            color = ORANGE_DARK if j == 0 else DARK_TEXT
            bold = j == 1
            add_text_box(slide, val, x, y, w, Inches(0.58),
                         font_size=11, color=color, bold=bold)


# ─────────────────────────────────────────────
# スライド10：AI搭載で実現できる機能
# ─────────────────────────────────────────────
def add_ai_features_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "△ AI搭載で実現できる機能",
                    "OpenAI APIキーを設定することで実現")

    features = [
        ("1", "個別化された次のステップアドバイス",
         "診断結果・支援者評価・日々の記録を総合分析\n「あなたの場合、まず○○を伸ばすといい」を生成"),
        ("2", "複数記録の時系列パターン認識",
         "感情ログ・タスク完了率の変化傾向を自動分析\n支援者向けAI所見レポートを自動生成"),
        ("3", "本人評価 × 支援者評価の統合分析",
         "複数の支援者コメントをAIが統合・要約\n矛盾する評価の意味を解釈し深い洞察を提供"),
        ("4", "職場シーンのコミュニケーション練習",
         "「遅刻しそうなとき、上司に何て言う？」\n職場でよくある場面をAIと繰り返し練習"),
        ("5", "アバターとの深い対話（完全版）",
         "現状はフォールバックで動作済み\nAPIキー設定で完全なAI対話・個別化された会話が実現"),
    ]

    y = Inches(1.6)
    for num, title, desc in features:
        item_h = Inches(1.02)
        add_card_shadow(slide, Inches(0.25), y, Inches(12.85), item_h)
        card = slide.shapes.add_shape(1, Inches(0.25), y, Inches(12.85), item_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        badge = slide.shapes.add_shape(1, Inches(0.25), y, Inches(0.45), item_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = ORANGE
        badge.line.fill.background()
        add_text_box(slide, num,
                     Inches(0.25), y + Inches(0.28),
                     Inches(0.45), Inches(0.45),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     Inches(0.85), y + Inches(0.04),
                     Inches(12.1), Inches(0.42),
                     font_size=15, bold=True, color=MID_GREEN)
        add_text_box(slide, desc,
                     Inches(0.85), y + Inches(0.46),
                     Inches(12.1), Inches(0.48),
                     font_size=12, color=DARK_TEXT)

        y += Inches(1.08)


# ─────────────────────────────────────────────
# スライド11：将来の壮大な構想
# ─────────────────────────────────────────────
def add_future_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "◇ 将来の構想",
                    "さらに大きなビジョンへ")

    visions = [
        ("🏢", "地元企業タイアップ・マッチング",
         "企業が「求める特性プロフィール」を登録\n利用者の診断結果と自動マッチング\n「あなたの特性にはこの企業が向いています」", MID_GREEN),
        ("📚", "就労実績データベース",
         "実際に就労した利用者の特性・就職先を匿名で蓄積\n「あなたと似た特性の方はここに就職しました」\n時間とともに精度が向上するデータ資産", PHASE2_BLUE),
        ("🤝", "多職種連携プラットフォーム",
         "作業療法士・就労支援員・ハローワーク\n担当者・企業が一つのプラットフォームで連携\n縦割り支援体制を横断的につなぐ", ORANGE_DARK),
        ("📱", "スマートフォンアプリ化",
         "プッシュ通知によるルーティン・タスクのリマインダー\nオフライン対応（インターネットなしでも動作）\n日常生活に溶け込んだ支援ツールへ", MID_GRAY),
    ]

    for i, (emoji, title, desc, accent) in enumerate(visions):
        col = i % 2
        row = i // 2
        x = Inches(0.25) + col * Inches(6.6)
        y = Inches(1.6) + row * Inches(2.65)
        w = Inches(6.3)
        h = Inches(2.45)

        add_card_shadow(slide, x, y, w, h)
        card = slide.shapes.add_shape(1, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        hdr_band = slide.shapes.add_shape(1, x, y, w, Inches(0.4))
        hdr_band.fill.solid()
        hdr_band.fill.fore_color.rgb = accent
        hdr_band.line.fill.background()

        add_text_box(slide, f"{emoji}  {title}",
                     x + Inches(0.15), y + Inches(0.05),
                     w - Inches(0.3), Inches(0.32),
                     font_size=14, bold=True, color=WHITE)

        add_text_box(slide, desc,
                     x + Inches(0.2), y + Inches(0.5),
                     w - Inches(0.4), Inches(1.8),
                     font_size=12, color=DARK_TEXT)


# ─────────────────────────────────────────────
# スライド12：全体ロードマップ
# ─────────────────────────────────────────────
def add_phases_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "開発ロードマップ",
                    "フェーズ別の全体像")

    phases = [
        ("Phase 1", "実装済み ★", MID_GREEN, [
            "利用者・支援者 登録/ログイン",
            "性格・特性診断（6特性・5タイプ）",
            "成長ロードマップ（3ステップ）",
            "日々の記録・ダッシュボード",
            "自分似のアバター機能",
            "AIチャット（GPT/フォールバック）",
            "支援者ダッシュボード",
            "漢字/ひらがな切り替え",
            "本番デプロイ（Railway）",
        ]),
        ("Phase 2", "AIなし追加 ◎", PHASE2_BLUE, [
            "支援者評価 × レーダーチャート",
            "希望と適性のギャップ表示",
            "進捗カレンダー + ストリーク",
            "朝・夕のルーティンチェック",
            "先輩事例データベース",
            "支援者間の連携ノート",
            "過去診断との比較グラフ",
        ]),
        ("Phase 3", "AI搭載 △", ORANGE_DARK, [
            "個別化ステップアドバイス",
            "時系列パターン認識",
            "本人×支援者評価の統合分析",
            "職場シーンのロールプレイ練習",
            "アバターとの深い対話（完全版）",
        ]),
        ("Phase 4", "将来構想 ◇", MID_GRAY, [
            "地元企業タイアップ・マッチング",
            "就労実績データベース",
            "多職種連携プラットフォーム",
            "スマートフォンアプリ化",
        ]),
    ]

    col_w = Inches(3.1)
    col_gap = Inches(0.14)
    x_start = Inches(0.25)
    y_top = Inches(1.55)
    hdr_h = Inches(0.85)
    body_h = Inches(5.1)

    for i, (phase, label, color, items) in enumerate(phases):
        x = x_start + i * (col_w + col_gap)

        add_card_shadow(slide, x, y_top, col_w, hdr_h + body_h)

        hdr = slide.shapes.add_shape(1, x, y_top, col_w, hdr_h)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()

        add_text_box(slide, phase,
                     x + Inches(0.1), y_top + Inches(0.04),
                     col_w - Inches(0.2), Inches(0.42),
                     font_size=16, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     x + Inches(0.1), y_top + Inches(0.46),
                     col_w - Inches(0.2), Inches(0.36),
                     font_size=11, color=WHITE,
                     align=PP_ALIGN.CENTER)

        if i < 3:
            arr = slide.shapes.add_shape(
                13,
                x + col_w + Inches(0.02), y_top + Inches(0.26),
                Inches(0.15), Inches(0.38)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = color
            arr.line.fill.background()

        body = slide.shapes.add_shape(1, x, y_top + hdr_h, col_w, body_h)
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        body.line.color.rgb = LIGHT_GREEN
        body.line.width = Pt(1)

        lb = slide.shapes.add_shape(1, x, y_top + hdr_h, Inches(0.06), body_h)
        lb.fill.solid()
        lb.fill.fore_color.rgb = color
        lb.line.fill.background()

        item_y = y_top + hdr_h + Inches(0.12)
        for item in items:
            add_text_box(slide, f"▶ {item}",
                         x + Inches(0.12), item_y,
                         col_w - Inches(0.18), Inches(0.5),
                         font_size=11, color=DARK_TEXT)
            item_y += Inches(0.52)


# ─────────────────────────────────────────────
# スライドA：AIをより賢くする仕組み：RAG
# ─────────────────────────────────────────────
def add_rag_intro_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "AIをより賢くする仕組み：RAG（ラグ）",
                    "RAG = Retrieval-Augmented Generation（検索して見つけてから答えるAI）")

    COLOR_NO_BG  = RGBColor(0xF5, 0xF5, 0xF5)
    COLOR_YES_BG = MINT

    # 左ボックス：普通のAI
    lx, ly, lw, lh = Inches(0.25), Inches(1.6), Inches(6.3), Inches(5.55)
    add_card_shadow(slide, lx, ly, lw, lh)
    left_box = slide.shapes.add_shape(1, lx, ly, lw, lh)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_NO_BG
    left_box.line.fill.background()

    hdr_l = slide.shapes.add_shape(1, lx, ly, lw, Inches(0.45))
    hdr_l.fill.solid()
    hdr_l.fill.fore_color.rgb = RGBColor(0x88, 0x44, 0x44)
    hdr_l.line.fill.background()

    add_text_box(slide, "❌  普通のAI（単体）",
                 lx + Inches(0.15), ly + Inches(0.06),
                 lw - Inches(0.3), Inches(0.36),
                 font_size=16, bold=True, color=WHITE)

    left_content = (
        "「清掃の仕事に向いてる人は？」\n"
        "\n"
        "↓\n"
        "\n"
        "一般的な知識だけで回答\n"
        "\n"
        "「几帳面で体力がある人が\n"
        "　向いています」\n"
        "\n"
        "→ あなた個人のことは知らない"
    )
    add_text_box(slide, left_content,
                 lx + Inches(0.2), ly + Inches(0.6),
                 lw - Inches(0.35), Inches(4.7),
                 font_size=14, color=DARK_TEXT)

    # 右ボックス：RAGあり
    rx, ry, rw, rh = Inches(6.8), Inches(1.6), Inches(6.3), Inches(5.55)
    add_card_shadow(slide, rx, ry, rw, rh)
    right_box = slide.shapes.add_shape(1, rx, ry, rw, rh)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_YES_BG
    right_box.line.fill.background()

    hdr_r = slide.shapes.add_shape(1, rx, ry, rw, Inches(0.45))
    hdr_r.fill.solid()
    hdr_r.fill.fore_color.rgb = MID_GREEN
    hdr_r.line.fill.background()

    add_text_box(slide, "✅  RAGあり",
                 rx + Inches(0.15), ry + Inches(0.06),
                 rw - Inches(0.3), Inches(0.36),
                 font_size=16, bold=True, color=WHITE)

    right_content = (
        "「清掃の仕事に向いてる人は？」\n"
        "\n"
        "↓　まずアプリのデータを検索\n"
        "　・田中さんの診断結果\n"
        "　・支援者の評価記録\n"
        "　・先輩の就職事例\n"
        "\n"
        "↓\n"
        "\n"
        "「田中さんは几帳面スコアが高く、\n"
        "先輩Aさんも同じ傾向で\n"
        "〇〇センターに就職しました」\n"
        "\n"
        "→ 個人に寄り添った回答"
    )
    add_text_box(slide, right_content,
                 rx + Inches(0.2), ry + Inches(0.6),
                 rw - Inches(0.35), Inches(4.7),
                 font_size=14, color=DARK_TEXT)

    add_text_box(slide, "VS",
                 Inches(6.32), Inches(4.15),
                 Inches(0.66), Inches(0.55),
                 font_size=18, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# スライドB：RAGでこのアプリができること
# ─────────────────────────────────────────────
def add_rag_usecase_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "RAGを使うと何が変わる？",
                    "アプリ内データを活用することで「個人に合った」回答が実現する")

    COLOR_NO_BG  = RGBColor(0xF5, 0xF5, 0xF5)

    cards = [
        (
            "💬  アバターとの会話",
            "❌ なし",
            "一般的な励ましの言葉",
            "✅ あり",
            "「先週より体調ログが改善してるね！\n今週も一緒にがんばろう」",
        ),
        (
            "🏢  就職先の提案",
            "❌ なし",
            "「清掃系が向いています」",
            "✅ あり",
            "「似た特性の山田さんは\n〇〇会社の清掃スタッフに\n就職しました」",
        ),
        (
            "📋  支援者向け所見",
            "❌ なし",
            "定型文のレポート",
            "✅ あり",
            "「3ヶ月の記録を見ると\n月曜に調子を崩しやすい\n傾向があります」",
        ),
    ]

    card_w = Inches(4.1)
    xs = [Inches(0.25), Inches(4.62), Inches(8.98)]

    for i, (card_title, no_label, no_text, yes_label, yes_text) in enumerate(cards):
        x = xs[i]
        y_top = Inches(1.55)

        add_card_shadow(slide, x, y_top, card_w, Inches(0.5))
        title_box = slide.shapes.add_shape(1, x, y_top, card_w, Inches(0.5))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = MID_GREEN
        title_box.line.fill.background()
        add_text_box(slide, card_title,
                     x + Inches(0.1), y_top + Inches(0.07),
                     card_w - Inches(0.2), Inches(0.38),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

        y_no = y_top + Inches(0.6)
        add_card_shadow(slide, x, y_no, card_w, Inches(2.0))
        no_box = slide.shapes.add_shape(1, x, y_no, card_w, Inches(2.0))
        no_box.fill.solid()
        no_box.fill.fore_color.rgb = COLOR_NO_BG
        no_box.line.color.rgb = SHADOW_GRAY
        no_box.line.width = Pt(1)
        add_text_box(slide, no_label,
                     x + Inches(0.1), y_no + Inches(0.08),
                     card_w - Inches(0.2), Inches(0.38),
                     font_size=13, bold=True, color=RGBColor(0x88, 0x44, 0x44))
        add_text_box(slide, no_text,
                     x + Inches(0.1), y_no + Inches(0.5),
                     card_w - Inches(0.2), Inches(1.3),
                     font_size=12, color=DARK_TEXT)

        y_yes = y_no + Inches(2.1)
        add_card_shadow(slide, x, y_yes, card_w, Inches(2.6))
        yes_box = slide.shapes.add_shape(1, x, y_yes, card_w, Inches(2.6))
        yes_box.fill.solid()
        yes_box.fill.fore_color.rgb = MINT
        yes_box.line.color.rgb = LIGHT_GREEN
        yes_box.line.width = Pt(1.5)
        add_text_box(slide, yes_label,
                     x + Inches(0.1), y_yes + Inches(0.08),
                     card_w - Inches(0.2), Inches(0.38),
                     font_size=13, bold=True, color=MID_GREEN)
        add_text_box(slide, yes_text,
                     x + Inches(0.1), y_yes + Inches(0.5),
                     card_w - Inches(0.2), Inches(1.9),
                     font_size=12, color=DARK_TEXT)


# ─────────────────────────────────────────────
# スライドC：このアプリのRAG構成イメージ
# ─────────────────────────────────────────────
def add_rag_architecture_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "データを活用するしくみ",
                    "RAG構成イメージ：利用者データがAIの「記憶」になる")

    sources = [
        "👤  利用者の診断結果・日々の記録",
        "👩‍💼  支援者の評価・観察メモ",
        "🏆  先輩の就職事例（匿名）",
        "🏢  地元企業の求人情報",
    ]

    src_w = Inches(5.85)
    src_xs = [Inches(0.25), Inches(7.23)]
    src_ys = [Inches(1.6), Inches(2.42)]

    for i, label in enumerate(sources):
        col = i % 2
        row = i // 2
        sx = src_xs[col]
        sy = src_ys[row]
        add_card_shadow(slide, sx, sy, src_w, Inches(0.62))
        box = slide.shapes.add_shape(1, sx, sy, src_w, Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LIGHT_GREEN
        box.line.width = Pt(1)

        lb = slide.shapes.add_shape(1, sx, sy, Inches(0.08), Inches(0.62))
        lb.fill.solid()
        lb.fill.fore_color.rgb = MID_GREEN
        lb.line.fill.background()

        add_text_box(slide, label,
                     sx + Inches(0.18), sy + Inches(0.1),
                     src_w - Inches(0.3), Inches(0.44),
                     font_size=13, color=DARK_TEXT)

    def add_down_arrow(slide, cx, y):
        arr = slide.shapes.add_shape(
            5,
            cx - Inches(0.2), y, Inches(0.4), Inches(0.35)
        )
        arr.fill.solid()
        arr.fill.fore_color.rgb = MID_GREEN
        arr.line.fill.background()

    add_down_arrow(slide, Inches(6.67), Inches(3.15))

    flow_items = [
        ("🗄️  データベース（蓄積・管理）", MID_GREEN, WHITE),
        ("🔍  質問に関連するデータを自動検索", PHASE2_BLUE, WHITE),
        ("🤖  AIが文脈を踏まえて回答", ORANGE_DARK, WHITE),
    ]

    flow_w = Inches(7.0)
    fx = Inches(3.17)
    fy = Inches(3.58)

    for j, (label, bg, fg) in enumerate(flow_items):
        add_card_shadow(slide, fx, fy, flow_w, Inches(0.62))
        box = slide.shapes.add_shape(1, fx, fy, flow_w, Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.fill.background()
        add_text_box(slide, label,
                     fx + Inches(0.2), fy + Inches(0.1),
                     flow_w - Inches(0.4), Inches(0.46),
                     font_size=15, bold=True, color=fg,
                     align=PP_ALIGN.CENTER)
        fy += Inches(0.62)
        if j < len(flow_items) - 1:
            add_down_arrow(slide, Inches(6.67), fy)
            fy += Inches(0.38)

    add_text_box(slide,
                 "💡 データが蓄積されるほどAIの精度が上がる「育てるAI」",
                 Inches(0.3), Inches(6.9),
                 Inches(12.7), Inches(0.42),
                 font_size=13, bold=True, color=MID_GREEN)


# ─────────────────────────────────────────────
# スライド16：フェーズ別 実装機能マップ
# ─────────────────────────────────────────────
def add_phase_feature_map_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "フェーズ別 実装機能マップ",
                    "現在から将来へ：段階的な機能拡張のロードマップ")

    COLOR_PHASE1_BG = RGBColor(0xD8, 0xF3, 0xDC)
    COLOR_PHASE2_BG = RGBColor(0xD1, 0xEC, 0xF1)
    COLOR_PHASE3_BG = RGBColor(0xFF, 0xE8, 0xD6)
    COLOR_PHASE4_BG = RGBColor(0xE9, 0xEC, 0xEF)

    phases = [
        (
            "Phase 1", "★ 実装済み", MID_GREEN, COLOR_PHASE1_BG,
            [
                "利用者・支援者登録/ログイン",
                "性格・特性診断（6特性）",
                "成長ロードマップ",
                "日々の記録・ダッシュボード",
                "自分似のアバター",
                "アバターとのAIチャット（基本）",
                "漢字/ひらがな切り替え",
                "支援者ダッシュボード",
                "Railway本番デプロイ済み",
            ]
        ),
        (
            "Phase 2", "◎ AIなしで追加可能", PHASE2_BLUE, COLOR_PHASE2_BG,
            [
                "支援者評価×多角的レーダーチャート",
                "希望と適性のギャップ表示",
                "次のステップ提示（ルールベース）",
                "進捗カレンダー＋継続ストリーク",
                "朝・夕のルーティンチェックリスト",
                "先輩事例データベース",
                "支援者間の連携ノート",
                "過去診断との比較グラフ",
            ]
        ),
        (
            "Phase 3", "△ AI搭載で強化", ORANGE_DARK, COLOR_PHASE3_BG,
            [
                "個別化された次のステップ（RAG）",
                "時系列パターン認識・AI所見",
                "本人×支援者評価の統合AI分析",
                "職場シーン練習（AIロールプレイ）",
                "アバターとの深い対話（完全AI）",
            ]
        ),
        (
            "Phase 4", "◇ 将来構想", MID_GRAY, COLOR_PHASE4_BG,
            [
                "地元企業タイアップ・マッチング",
                "就労実績データベース（先輩事例）",
                "多職種連携プラットフォーム",
                "スマートフォンアプリ化",
            ]
        ),
    ]

    col_w = Inches(3.1)
    col_gap = Inches(0.13)
    x_start = Inches(0.25)
    y_top = Inches(1.55)
    hdr_h = Inches(0.82)
    body_h = Inches(5.1)

    for i, (phase, label, hdr_color, bg_color, items) in enumerate(phases):
        x = x_start + i * (col_w + col_gap)

        add_card_shadow(slide, x, y_top, col_w, hdr_h + body_h)

        hdr = slide.shapes.add_shape(1, x, y_top, col_w, hdr_h)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = hdr_color
        hdr.line.fill.background()

        add_text_box(slide, phase,
                     x + Inches(0.1), y_top + Inches(0.04),
                     col_w - Inches(0.2), Inches(0.38),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     x + Inches(0.1), y_top + Inches(0.43),
                     col_w - Inches(0.2), Inches(0.34),
                     font_size=10, color=WHITE,
                     align=PP_ALIGN.CENTER)

        body = slide.shapes.add_shape(1, x, y_top + hdr_h, col_w, body_h)
        body.fill.solid()
        body.fill.fore_color.rgb = bg_color
        body.line.color.rgb = hdr_color
        body.line.width = Pt(1)

        item_y = y_top + hdr_h + Inches(0.12)
        for item in items:
            add_text_box(slide, f"▶ {item}",
                         x + Inches(0.12), item_y,
                         col_w - Inches(0.2), Inches(0.47),
                         font_size=10, color=DARK_TEXT)
            item_y += Inches(0.5)


# ─────────────────────────────────────────────
# スライド17：フェーズ別 費用シミュレーション
# ─────────────────────────────────────────────
def add_phase_cost_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "フェーズ別 月額費用の目安",
                    "利用者規模に応じた従量課金制・小規模なら実質ほぼ無料")

    COLOR_PHASE1_BG = RGBColor(0xD8, 0xF3, 0xDC)
    COLOR_PHASE2_BG = RGBColor(0xD1, 0xEC, 0xF1)
    COLOR_PHASE3_BG = RGBColor(0xFF, 0xE8, 0xD6)
    COLOR_PHASE4_BG = RGBColor(0xE9, 0xEC, 0xEF)

    headers = ["フェーズ", "主な追加機能", "利用者規模", "Railway", "OpenAI API", "月額合計"]
    col_xs = [Inches(0.25), Inches(2.1), Inches(5.2), Inches(6.9), Inches(8.6), Inches(10.3)]
    col_ws = [Inches(1.8), Inches(3.0), Inches(1.6), Inches(1.6), Inches(1.65), Inches(2.8)]

    hdr_y = Inches(1.6)
    hdr_h = Inches(0.48)

    hdr_bg = slide.shapes.add_shape(1, Inches(0.25), hdr_y, Inches(12.85), hdr_h)
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = MID_GREEN
    hdr_bg.line.fill.background()

    for h, x, w in zip(headers, col_xs, col_ws):
        add_text_box(slide, h, x + Inches(0.05), hdr_y + Inches(0.06),
                     w - Inches(0.1), Inches(0.38),
                     font_size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

    rows = [
        ("Phase 1\n（現状）",
         "基本機能・AIチャット",
         "〜30人",
         "$0〜5",
         "$1〜3",
         "約$1〜8\n（〜1,200円）",
         COLOR_PHASE1_BG),
        ("Phase 2\n（AIなし追加）",
         "評価・記録・事例DB",
         "30〜100人",
         "$5〜20",
         "$3〜10",
         "約$8〜30\n（〜4,500円）",
         COLOR_PHASE2_BG),
        ("Phase 3\n（AI強化）",
         "RAG・個別分析",
         "50〜200人",
         "$20",
         "$20〜60",
         "約$40〜80\n（〜12,000円）",
         COLOR_PHASE3_BG),
        ("Phase 4\n（将来構想）",
         "企業連携・アプリ化",
         "200人以上",
         "$20〜50",
         "$50〜150",
         "約$70〜200\n（〜30,000円）",
         COLOR_PHASE4_BG),
    ]

    row_h = Inches(1.0)
    row_y = hdr_y + hdr_h + Inches(0.05)

    for row_data in rows:
        *cells, bg_color = row_data
        row_bg = slide.shapes.add_shape(1, Inches(0.25), row_y, Inches(12.85), row_h)
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.color.rgb = LIGHT_GREEN
        row_bg.line.width = Pt(0.5)

        for j, (val, x, w) in enumerate(zip(cells, col_xs, col_ws)):
            bold = j == 5
            color = MID_GREEN if j == 5 else DARK_TEXT
            add_text_box(slide, val, x + Inches(0.05), row_y + Inches(0.1),
                         w - Inches(0.1), row_h - Inches(0.15),
                         font_size=11, bold=bold, color=color,
                         align=PP_ALIGN.CENTER)

        row_y += row_h + Inches(0.04)

    notes = [
        "💡 OpenAI APIはGPT-4o-miniを使用。1チャット約0.1円以下の低コスト",
        "💡 利用者が少ない間はほぼ無料。使った分だけ払う従量課金制",
        "💡 Railwayは月$5のクレジット付与あり（小規模なら実質無料）",
    ]
    note_y = row_y + Inches(0.04)
    for note in notes:
        add_text_box(slide, note,
                     Inches(0.3), note_y,
                     Inches(12.7), Inches(0.3),
                     font_size=10, color=MID_GRAY, italic=True)
        note_y += Inches(0.28)


# ─────────────────────────────────────────────
# スライド18：まとめ・次のアクション
# ─────────────────────────────────────────────
def add_summary_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg_color(slide, OFF_WHITE)
    add_slide_decorations(slide)
    add_slide_title(slide, "次のステップを一緒に決めましょう",
                    "ご確認・ご検討いただきたい項目")

    # 左カラム：検討事項
    add_text_box(slide, "📋 検討事項",
                 Inches(0.3), Inches(1.6),
                 Inches(6.5), Inches(0.45),
                 font_size=18, bold=True, color=DARK_TEXT)

    checks = [
        "Phase 2（AIなし追加機能）の優先順位を決める",
        "OpenAI APIキーの導入・費用の検討",
        "本番環境での利用者テスト実施",
        "支援者・利用者へのフィードバック収集",
        "Phase 3以降のAI機能実装スケジュール",
        "地元企業・支援機関との連携検討（Phase 4）",
    ]

    y = Inches(2.15)
    for item in checks:
        item_h = Inches(0.68)
        add_card_shadow(slide, Inches(0.3), y, Inches(6.4), item_h)
        card = slide.shapes.add_shape(1, Inches(0.3), y, Inches(6.4), item_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()

        chk = slide.shapes.add_shape(1, Inches(0.3), y, Inches(0.1), item_h)
        chk.fill.solid()
        chk.fill.fore_color.rgb = MID_GREEN
        chk.line.fill.background()

        dot = slide.shapes.add_shape(9, Inches(0.5), y + Inches(0.2), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE
        dot.line.fill.background()

        add_text_box(slide, item,
                     Inches(0.9), y + Inches(0.1),
                     Inches(5.7), Inches(0.5),
                     font_size=14, color=DARK_TEXT)
        y += Inches(0.76)

    # 右カラム：URL情報
    rx, ry, rw, rh = Inches(7.1), Inches(1.55), Inches(6.0), Inches(5.6)
    add_card_shadow(slide, rx, ry, rw, rh)
    info_box = slide.shapes.add_shape(1, rx, ry, rw, rh)
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = DARK_GREEN
    info_box.line.fill.background()

    add_text_box(slide, "🚀 現在の本番URL",
                 rx + Inches(0.2), ry + Inches(0.25),
                 rw - Inches(0.4), Inches(0.5),
                 font_size=17, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "https://web-production-c0abe.up.railway.app/",
                 rx + Inches(0.2), ry + Inches(0.85),
                 rw - Inches(0.4), Inches(0.5),
                 font_size=11, color=LIGHT_GREEN,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "すでに動いています！\n今すぐお試しいただけます",
                 rx + Inches(0.2), ry + Inches(1.45),
                 rw - Inches(0.4), Inches(0.8),
                 font_size=17, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    sep = slide.shapes.add_shape(1, rx + Inches(0.3), ry + Inches(2.38),
                                 rw - Inches(0.6), Inches(0.04))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ORANGE
    sep.line.fill.background()

    add_text_box(slide, "ご質問・ご要望はお気軽にどうぞ",
                 rx + Inches(0.2), ry + Inches(2.52),
                 rw - Inches(0.4), Inches(0.45),
                 font_size=14, color=LIGHT_GREEN,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "一緒に「ステップアップナビ」を\n育てていきましょう 🌱",
                 rx + Inches(0.2), ry + Inches(3.1),
                 rw - Inches(0.4), Inches(1.3),
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────
if __name__ == '__main__':
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("スライドを作成中...")

    add_title_slide(prs)
    print("  [1/18] 表紙")

    add_overview_slide(prs)
    print("  [2/18] アプリ概要")

    add_current_features_overview(prs)
    print("  [3/18] 現在の実装機能（概要）")

    add_diagnosis_slide(prs)
    print("  [4/18] 性格・特性診断")

    add_roadmap_slide(prs)
    print("  [5/18] 成長ロードマップ")

    add_avatar_slide(prs)
    print("  [6/18] アバター機能")

    add_ai_chat_slide(prs)
    print("  [7/18] AIチャット")

    add_accessibility_slide(prs)
    print("  [8/18] 漢字切り替え・支援者ダッシュボード")

    add_no_ai_features_slide(prs)
    print("  [9/18] AIなしで追加できる機能")

    add_ai_features_slide(prs)
    print("  [10/18] AI搭載で実現できる機能")

    add_future_slide(prs)
    print("  [11/18] 将来の構想")

    add_phases_slide(prs)
    print("  [12/18] 全体ロードマップ")

    add_rag_intro_slide(prs)
    print("  [13/18] RAG紹介：AIをより賢くする仕組み")

    add_rag_usecase_slide(prs)
    print("  [14/18] RAG活用：何が変わる？")

    add_rag_architecture_slide(prs)
    print("  [15/18] RAG構成イメージ")

    add_phase_feature_map_slide(prs)
    print("  [16/18] フェーズ別 実装機能マップ")

    add_phase_cost_slide(prs)
    print("  [17/18] フェーズ別 費用シミュレーション")

    add_summary_slide(prs)
    print("  [18/18] まとめ・次のアクション")

    output_path = '/Users/matsunaganaoto/Desktop/projects/App/ステップアップナビ_機能紹介.pptx'
    prs.save(output_path)
    print(f"\n✅ 完成しました！")
    print(f"📄 ファイル: {output_path}")
    print(f"📊 スライド数: {len(prs.slides)} 枚")
