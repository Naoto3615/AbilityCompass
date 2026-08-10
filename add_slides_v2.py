"""
東濃信用金庫様_提案資料.pptx に3枚のスライドを追加してv2として保存するスクリプト
挿入位置: スライド8（開発スケジュール）の後、スライド9（300万円）の前
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# カラー定義
BG_COLOR = RGBColor(0x0f, 0x17, 0x2a)       # ダークネイビー
GREEN    = RGBColor(0x10, 0xb9, 0x81)        # エメラルドグリーン
BLUE     = RGBColor(0x3b, 0x82, 0xf6)        # ブルー
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)        # アンバー
WHITE    = RGBColor(0xff, 0xff, 0xff)
GRAY     = RGBColor(0x94, 0xa3, 0xb8)        # スレートグレー
DARK_CARD= RGBColor(0x1e, 0x29, 0x3b)        # カードBackground

INPUT  = '/Users/matsunaganaoto/Desktop/projects/App/東濃信用金庫様_提案資料.pptx'
OUTPUT = '/Users/matsunaganaoto/Desktop/projects/App/東濃信用金庫様_提案資料_v2.pptx'

prs = Presentation(INPUT)

# スライドサイズ確認
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
print(f"Slide size: {SLIDE_W.inches:.2f}\" x {SLIDE_H.inches:.2f}\"")

# --- ヘルパー関数 ---

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, x, y, w, h, fill_color, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if radius:
        sp = shape._element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is None:
            pass
        else:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                for child in list(avLst):
                    avLst.remove(child)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False,
                line_spacing=None):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if line_spacing:
        from pptx.util import Pt as _Pt
        p.line_spacing = line_spacing
    return txBox

def add_multiline_textbox(slide, lines, x, y, w, h,
                          font_size=11, bold_first=False,
                          color=WHITE, first_color=None,
                          first_size=None, line_color=None):
    """lines: list of str. 最初の行をbold/別色にする場合はbold_first=True"""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(first_size if (i == 0 and first_size) else font_size)
        run.font.bold = (bold_first and i == 0)
        c = first_color if (i == 0 and first_color) else color
        if line_color and i > 0:
            c = line_color
        run.font.color.rgb = c
    return txBox

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)

def add_slide_title(slide, title, subtitle=None):
    """スライド共通タイトルブロック"""
    # 上部のカラーバー
    bar = add_rect(slide, 0, 0, SLIDE_W.inches, 0.08, GREEN)
    
    # タイトルテキスト
    tb = add_textbox(slide, title,
                     x=0.4, y=0.15, w=SLIDE_W.inches - 0.8, h=0.55,
                     font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    x=0.4, y=0.68, w=SLIDE_W.inches - 0.8, h=0.35,
                    font_size=13, color=GRAY)

def add_card(slide, x, y, w, h, accent_color, title, body_lines,
             title_size=13, body_size=10.5, corner_radius=20000):
    """カードコンポーネント: 背景+アクセントライン+タイトル+本文"""
    # カード背景
    card = add_rect(slide, x, y, w, h, DARK_CARD, radius=corner_radius)
    # アクセントボーダー（左）
    accent = add_rect(slide, x, y, 0.04, h, accent_color)
    # タイトル
    add_textbox(slide, title,
                x=x+0.12, y=y+0.08, w=w-0.2, h=0.35,
                font_size=title_size, bold=True, color=accent_color)
    # 本文
    add_multiline_textbox(slide, body_lines,
                          x=x+0.12, y=y+0.42, w=w-0.2, h=h-0.55,
                          font_size=body_size, color=WHITE)

def add_bullet_textbox(slide, items, x, y, w, h, bullet='•',
                       font_size=10.5, color=WHITE, indent=0.15):
    """箇条書きテキストボックス"""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet} {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox

# ================================================================
# スライドA: 市場規模・需要の根拠
# ================================================================
def create_slide_a(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG_COLOR)
    
    add_slide_title(
        slide,
        "解決すべき課題の大きさ",
        "〜市場データが示すニーズ〜"
    )
    
    W = SLIDE_W.inches
    H = SLIDE_H.inches
    
    # === 左カード: 就労支援市場 ===
    card_x = 0.35
    card_y = 1.1
    card_w = (W - 0.9) / 2
    card_h = H - card_y - 0.85
    
    add_rect(slide, card_x, card_y, card_w, card_h, DARK_CARD, radius=25000)
    add_rect(slide, card_x, card_y, 0.05, card_h, GREEN)
    
    add_textbox(slide, "就労支援市場",
                x=card_x+0.14, y=card_y+0.1, w=card_w-0.2, h=0.35,
                font_size=14, bold=True, color=GREEN)
    
    left_items = [
        "全国の障がい者雇用者数：約64万人（2023年、厚労省）",
        "法定雇用率（2.3%）未達成企業：約47%（中小企業中心）",
        "障がい者の離職率：約30%が1年以内に離職（定着支援の需要大）",
        "東濃地区の就労継続支援事業所：約40事業所以上",
    ]
    add_bullet_textbox(slide, left_items,
                       x=card_x+0.14, y=card_y+0.52, w=card_w-0.25,
                       h=card_h-0.65, font_size=11, color=WHITE)
    
    # === 右カード: 放課後デイサービス市場 ===
    card2_x = card_x + card_w + 0.2
    add_rect(slide, card2_x, card_y, card_w, card_h, DARK_CARD, radius=25000)
    add_rect(slide, card2_x, card_y, 0.05, card_h, BLUE)
    
    add_textbox(slide, "放課後デイサービス市場",
                x=card2_x+0.14, y=card_y+0.1, w=card_w-0.2, h=0.35,
                font_size=14, bold=True, color=BLUE)
    
    right_items = [
        "全国の放課後デイサービス事業所数：約17,000箇所（年々増加中）",
        "利用者数：約30万人（2022年度、厚労省）",
        "支援員の離職・バーンアウトが業界課題",
        "岐阜県内の放課後デイ事業所：約400箇所以上",
    ]
    add_bullet_textbox(slide, right_items,
                       x=card2_x+0.14, y=card_y+0.52, w=card_w-0.25,
                       h=card_h-0.65, font_size=11, color=WHITE)
    
    # 補足
    add_textbox(slide,
                "※出典：厚生労働省「障害者雇用状況の集計結果」「障害福祉サービス等の利用状況について」",
                x=0.35, y=H-0.65, w=W-0.7, h=0.35,
                font_size=8.5, color=GRAY)
    
    return slide

# ================================================================
# スライドB: 収益モデル
# ================================================================
def create_slide_b(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG_COLOR)
    
    add_slide_title(
        slide,
        "クラファン後も続く、持続可能なビジネスモデル",
        ""
    )
    
    W = SLIDE_W.inches
    H = SLIDE_H.inches
    
    col_w = (W - 0.6) / 3
    card_y = 1.05
    card_h = H - card_y - 0.95
    gap = 0.1
    
    cards = [
        {
            "x": 0.3,
            "color": GREEN,
            "label": "就労者・個人向け",
            "items": [
                "月額サブスク（個人利用）",
                "  月500〜1,000円",
                "",
                "就労支援事業所向けライセンス",
                "  月5,000〜10,000円/事業所",
            ]
        },
        {
            "x": 0.3 + col_w + gap,
            "color": BLUE,
            "label": "企業・事業所向け",
            "items": [
                "採用企業向けダッシュボード",
                "  月10,000〜30,000円/社",
                "",
                "放課後デイ事業所向けパック",
                "  月8,000〜15,000円/事業所",
                "",
                "初期導入支援費",
                "  50,000〜100,000円",
            ]
        },
        {
            "x": 0.3 + (col_w + gap) * 2,
            "color": AMBER,
            "label": "行政・補助金連携",
            "items": [
                "障害者就労支援関連の",
                "補助金・助成金活用",
                "",
                "自治体向けSaaS提供",
                "（ライセンス販売）",
                "",
                "社会的インパクト投資・",
                "ESG資金の活用",
            ]
        },
    ]
    
    for c in cards:
        add_rect(slide, c["x"], card_y, col_w, card_h, DARK_CARD, radius=25000)
        add_rect(slide, c["x"], card_y, 0.05, card_h, c["color"])
        
        # アイコン的な丸+番号
        add_textbox(slide, c["label"],
                    x=c["x"]+0.14, y=card_y+0.1, w=col_w-0.2, h=0.4,
                    font_size=13, bold=True, color=c["color"])
        
        # 仕切り線
        line_box = add_rect(slide, c["x"]+0.14, card_y+0.52, col_w-0.28, 0.01,
                            RGBColor(0x33, 0x44, 0x55))
        
        add_multiline_textbox(slide, c["items"],
                              x=c["x"]+0.14, y=card_y+0.6, w=col_w-0.2,
                              h=card_h-0.72, font_size=10.5, color=WHITE)
    
    # 目標テキスト
    goal_y = H - 0.85
    add_rect(slide, 0.3, goal_y, W-0.6, 0.45, RGBColor(0x16, 0x1f, 0x30), radius=15000)
    add_textbox(slide,
                "目標：リリース1年以内に月額収益30万円超・単月黒字化を目指す",
                x=0.5, y=goal_y+0.05, w=W-1.0, h=0.35,
                font_size=11.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    
    return slide

# ================================================================
# スライドC: 競合との差別化
# ================================================================
def create_slide_c(prs):
    slide = add_blank_slide(prs)
    set_bg(slide, BG_COLOR)
    
    add_slide_title(
        slide,
        "既存サービスにはない、3つの独自価値",
        ""
    )
    
    W = SLIDE_W.inches
    H = SLIDE_H.inches
    
    col_w = (W - 0.6) / 3
    card_y = 1.05
    card_h = H - card_y - 0.5
    gap = 0.1
    
    cards = [
        {
            "x": 0.3,
            "color": GREEN,
            "num": "01",
            "label": "当事者家族が作るUX",
            "subtitle": "使う人の気持ちが分かる設計",
            "body": [
                "開発者自身が障がい児の父親。",
                "既存ツールの「難しくて使えない」",
                "「現場感がない」という声に応え、",
                "徹底した使いやすさを追求。",
                "",
                "支援員・利用者・家族の三者が",
                "自然に使えるUIを実現。",
            ]
        },
        {
            "x": 0.3 + col_w + gap,
            "color": BLUE,
            "num": "02",
            "label": "就労支援×放課後デイの一体化",
            "subtitle": "子どもから大人まで、切れ目のない支援",
            "body": [
                "就労支援と放課後デイを一つの",
                "プラットフォームで提供するサービスは",
                "国内でも希少。",
                "",
                "子ども時代の発達記録を大人の就労",
                "支援につなげる「ライフステージを",
                "超えた継続支援」が実現できる。",
            ]
        },
        {
            "x": 0.3 + (col_w + gap) * 2,
            "color": AMBER,
            "num": "03",
            "label": "AIアバターによる24時間サポート",
            "subtitle": "支援員がいない時間も、孤立させない",
            "body": [
                "自分に似たアバターとのAIチャットで、",
                "悩みや不安をいつでも吐き出せる。",
                "",
                "従来の支援ツールにはない",
                "「感情的なサポート機能」が、",
                "就労定着率向上に直結する。",
            ]
        },
    ]
    
    for c in cards:
        add_rect(slide, c["x"], card_y, col_w, card_h, DARK_CARD, radius=25000)
        add_rect(slide, c["x"], card_y, col_w, 0.06, c["color"])  # 上部カラーバー
        
        # 番号
        add_textbox(slide, c["num"],
                    x=c["x"]+0.15, y=card_y+0.1, w=0.7, h=0.45,
                    font_size=28, bold=True, color=c["color"])
        
        # ラベル（カテゴリ名）
        add_textbox(slide, c["label"],
                    x=c["x"]+0.15, y=card_y+0.55, w=col_w-0.25, h=0.4,
                    font_size=12, bold=True, color=c["color"])
        
        # サブタイトル
        add_textbox(slide, c["subtitle"],
                    x=c["x"]+0.15, y=card_y+0.92, w=col_w-0.25, h=0.38,
                    font_size=10.5, bold=False, color=WHITE)
        
        # 仕切り
        add_rect(slide, c["x"]+0.15, card_y+1.3, col_w-0.3, 0.01,
                 RGBColor(0x33, 0x44, 0x55))
        
        # 本文
        add_multiline_textbox(slide, c["body"],
                              x=c["x"]+0.15, y=card_y+1.38, w=col_w-0.25,
                              h=card_h-1.52, font_size=10, color=GRAY)
    
    return slide

# ================================================================
# メイン処理: スライドを末尾に追加してXMLで順序を変更
# ================================================================

print(f"Before: {len(prs.slides)} slides")

# 3枚を末尾に追加（インデックス10, 11, 12）
slide_a = create_slide_a(prs)
slide_b = create_slide_b(prs)
slide_c = create_slide_c(prs)

print(f"After adding: {len(prs.slides)} slides")

# --- スライド順序の並び替え ---
# 現在: 0-9(既存10枚) + 10,11,12(新3枚)
# 目標: 0-7(既存1-8), 10,11,12(新A,B,C), 8,9(既存9,10)

def move_slide(prs, old_index, new_index):
    """スライドをold_indexからnew_indexへ移動"""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    slide_to_move = slides_list[old_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(new_index, slide_to_move)

# 現在の順序（追加後）:
# 0,1,2,3,4,5,6,7 = 既存1-8
# 8,9             = 既存9-10
# 10,11,12        = 新A,B,C
#
# 目標:
# 0,1,2,3,4,5,6,7 = 既存1-8 (そのまま)
# 8,9,10          = 新A,B,C (現在の10,11,12)
# 11,12           = 既存9-10 (現在の8,9)
#
# 操作:
# 1. インデックス10(新A)を8に移動
# 2. インデックス11(新B、移動後)を9に移動  
# 3. インデックス12(新C、移動後)を10に移動

move_slide(prs, 10, 8)   # 新A: 10→8
move_slide(prs, 11, 9)   # 新B: 11→9
move_slide(prs, 12, 10)  # 新C: 12→10

print(f"Final: {len(prs.slides)} slides")

# 保存
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")

# 確認
prs2 = Presentation(OUTPUT)
print(f"\n=== 最終スライド一覧 ===")
for i, slide in enumerate(prs2.slides):
    first_text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            first_text = shape.text.strip()[:50]
            break
    print(f"Slide {i+1}: {first_text}")
